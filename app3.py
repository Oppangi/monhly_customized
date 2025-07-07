import streamlit as st
import os
import tempfile
import json
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional, Union
import zipfile
import io
import base64
import re
import boto3
from botocore.exceptions import ClientError
import PyPDF2
from langdetect import detect
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from dotenv import load_dotenv
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import numpy as np
from PIL import Image
import urllib.request

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(), logging.FileHandler("app.log")]
)
logger = logging.getLogger(__name__)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        padding: 20px;
        border-radius: 10px;
        color: white;
        margin-bottom: 30px;
        text-align: center;
    }
    .feature-box {
        background: #f8f9fa;
        padding: 15px;
        border-radius: 8px;
        border-left: 4px solid #667eea;
        margin: 10px 0;
    }
    .metric-container {
        background: white;
        padding: 20px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin: 10px 0;
    }
    .preview-image {
        max-width: 100%;
        border-radius: 8px;
        margin-top: 10px;
    }
    .template-option {
        padding: 15px;
        border: 1px solid #ddd;
        border-radius: 8px;
        margin-bottom: 15px;
    }
</style>
""", unsafe_allow_html=True)

class Settings:
    def __init__(self):
        load_dotenv()
        self.aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID')
        self.aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY')
        self.aws_session_token = os.getenv('AWS_SESSION_TOKEN')
        self.aws_default_region = os.getenv('AWS_DEFAULT_REGION', 'us-east-1')
        self.bedrock_model_id = os.getenv('BEDROCK_MODEL_ID', 'anthropic.claude-3-sonnet-20240229-v1:0')
        self.bedrock_region = os.getenv('BEDROCK_REGION', 'us-east-1')
       
        # Template handling
        self.template_path = None
        self.use_custom_template = False
        
        self.upload_folder = Path('uploads')
        self.output_folder = Path('output')
        self.template_folder = Path('user_templates')
        self.max_file_size = 10485760  # 10MB
        self.allowed_extensions = ['pdf', 'pptx']
        self.default_source_language = 'auto'
        self.default_target_language = 'en'
        self.supported_languages = ['en', 'es', 'fr', 'de', 'it', 'pt', 'nl', 'ru', 'zh', 'ja', 'ko', 'ar', 'hi']
       
        # Create directories if they don't exist
        self.upload_folder.mkdir(exist_ok=True)
        self.output_folder.mkdir(exist_ok=True)
        self.template_folder.mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates').mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates' / 'styles').mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates' / 'icons').mkdir(exist_ok=True)

        # Set up AWS session
        try:
            boto3.setup_default_session(
                aws_access_key_id=self.aws_access_key_id,
                aws_secret_access_key=self.aws_secret_access_key,
                aws_session_token=self.aws_session_token,
                region_name=self.aws_default_region
            )
        except Exception as e:
            st.warning(f"AWS setup warning: {str(e)}")
            logger.error(f"AWS setup error: {str(e)}")
    
    def set_template_path(self, template_path: Union[str, Path]):
        """Set the template path based on user input"""
        if not template_path:
            self.template_path = None
            self.use_custom_template = False
            return
            
        try:
            if isinstance(template_path, str):
                if template_path.startswith(('http://', 'https://')):
                    # Download from URL
                    template_name = os.path.basename(template_path)
                    local_path = self.template_folder / template_name
                    urllib.request.urlretrieve(template_path, local_path)
                    self.template_path = local_path
                else:
                    self.template_path = Path(template_path)
            else:
                self.template_path = template_path
                
            self.use_custom_template = True
            logger.info(f"Using custom template: {self.template_path}")
        except Exception as e:
            st.error(f"Error setting template path: {str(e)}")
            logger.error(f"Template path error: {str(e)}")
            self.template_path = None
            self.use_custom_template = False

class DocumentProcessor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
   
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        try:
            text = ""
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return self._clean_text(text)
        except Exception as e:
            self.logger.error(f"Error extracting text from PDF: {str(e)}")
            st.error(f"Error extracting PDF: {str(e)}")
            return ""
   
    def process_text_document(self, content: str) -> str:
        return self._clean_text(content)
   
    def _clean_text(self, text: str) -> str:
        text = ' '.join(text.split())
        text = text.replace('\x00', '').replace('\ufeff', '')
        return text
   
    def detect_language(self, text: str) -> Optional[str]:
        try:
            if len(text.strip()) < 10:
                return None
            return detect(text)
        except:
            return None
   
    def extract_tables(self, text: str) -> List[Dict[str, List]]:
        """Extract table-like data from text"""
        tables = []
        lines = text.split('\n')
        current_table = None
        headers = None
       
        for line in lines:
            if re.search(r'[,|]\s*|\d+\s+\w+\s+\d+', line):
                items = [item.strip() for item in re.split(r'[,|]\s*|\s{2,}', line) if item.strip()]
                if len(items) >= 2:
                    if not headers:
                        headers = items
                        current_table = {"headers": headers, "rows": []}
                    else:
                        if len(items) == len(headers):
                            current_table["rows"].append(items)
            else:
                if current_table:
                    tables.append(current_table)
                    headers = None
                    current_table = None
       
        if current_table:
            tables.append(current_table)
       
        return tables
   
    def extract_statistical_data(self, text: str) -> Dict[str, Any]:
        """Extract numerical data, percentages, and statistical information from text"""
        statistical_data = {
            'numbers': [],
            'percentages': [],
            'dates': [],
            'monetary_values': [],
            'metrics': [],
            'comparisons': []
        }
       
        # Extract numbers with context
        number_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(%|percent|million|billion|thousand|units|kg|tons|miles|km)?'
        numbers = re.findall(number_pattern, text, re.IGNORECASE)
       
        for context, value, unit in numbers:
            try:
                numeric_value = float(value.replace(',', ''))
                statistical_data['numbers'].append({
                    'context': context.strip(),
                    'value': numeric_value,
                    'unit': unit,
                    'original': f"{value} {unit}".strip()
                })
            except ValueError:
                continue
       
        # Extract percentages
        percentage_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:\.\d+)?)\s*(%|percent)'
        percentages = re.findall(percentage_pattern, text, re.IGNORECASE)
       
        for context, value, unit in percentages:
            try:
                statistical_data['percentages'].append({
                    'context': context.strip(),
                    'value': float(value),
                    'original': f"{value}%"
                })
            except ValueError:
                continue
       
        # Extract monetary values
        money_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*\$(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|thousand|M|B|K)?'
        monetary = re.findall(money_pattern, text, re.IGNORECASE)
       
        for context, value, multiplier in monetary:
            try:
                numeric_value = float(value.replace(',', ''))
                multiplier_map = {'million': 1000000, 'M': 1000000, 'billion': 1000000000, 'B': 1000000000, 'thousand': 1000, 'K': 1000}
                if multiplier and multiplier in multiplier_map:
                    numeric_value *= multiplier_map[multiplier]
               
                statistical_data['monetary_values'].append({
                    'context': context.strip(),
                    'value': numeric_value,
                    'original': f"${value} {multiplier}".strip()
                })
            except ValueError:
                continue
       
        # Extract comparison data (increased by X%, decreased by Y%)
        comparison_pattern = r'(\w+(?:\s+\w+){0,2})\s*(?:increased|decreased|rose|fell|grew|declined)\s*(?:by\s*)?(\d+(?:\.\d+)?)\s*(%|percent)'
        comparisons = re.findall(comparison_pattern, text, re.IGNORECASE)
       
        for context, value, unit in comparisons:
            try:
                statistical_data['comparisons'].append({
                    'context': context.strip(),
                    'value': float(value),
                    'type': 'increase' if any(word in text.lower() for word in ['increased', 'rose', 'grew']) else 'decrease'
                })
            except ValueError:
                continue
       
        return statistical_data
   
    def analyze_content_type(self, text: str) -> Dict[str, Any]:
        content_analysis = {
            'type': 'general',
            'has_numbers': bool(re.search(r'\d+', text)),
            'has_dates': bool(re.search(r'\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b', text, re.I)),
            'has_tables': bool(self.extract_tables(text)),
            'topics': [],
            'potential_charts': []
        }
        text_lower = text.lower()
        logistics_keywords = ['logistics', 'cargo', 'fleet', 'vehicle', 'transport', 'shipping', 'warehouse']
        business_keywords = ['revenue', 'profit', 'sales', 'market', 'customer', 'performance']
        story_keywords = ['story', 'narrative', 'character', 'plot', 'once upon']
        sop_keywords = ['procedure', 'sop', 'standard operating', 'guideline', 'protocol', 'instruction']
       
        if any(keyword in text_lower for keyword in logistics_keywords):
            content_analysis['type'] = 'logistics'
            content_analysis['topics'] = ['Fleet Operations', 'Cargo Management', 'Delivery Metrics']
        elif any(keyword in text_lower for keyword in business_keywords):
            content_analysis['type'] = 'business'
            content_analysis['topics'] = ['Financials', 'Market Trends', 'Performance']
        elif any(keyword in text_lower for keyword in story_keywords):
            content_analysis['type'] = 'story'
            content_analysis['topics'] = ['Narrative', 'Themes', 'Characters']
        elif any(keyword in text_lower for keyword in sop_keywords):
            content_analysis['type'] = 'sop'
            content_analysis['topics'] = ['Procedures', 'Guidelines', 'Compliance']
       
        if content_analysis['has_numbers'] or content_analysis['has_tables']:
            content_analysis['potential_charts'] = ['bar', 'line', 'pie', 'table']
       
        # Add statistical analysis
        statistical_data = self.extract_statistical_data(text)
        content_analysis['statistical_data'] = statistical_data
        content_analysis['has_statistics'] = len(statistical_data['numbers']) > 0 or len(statistical_data['percentages']) > 0

        # Suggest chart types based on data
        if statistical_data['percentages']:
            content_analysis['potential_charts'].extend(['pie', 'donut'])
        if statistical_data['comparisons']:
            content_analysis['potential_charts'].extend(['bar', 'column'])
        if statistical_data['monetary_values']:
            content_analysis['potential_charts'].extend(['line', 'area'])
       
        return content_analysis

class BedrockClient:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.logger = logging.getLogger(__name__)
        try:
            self.bedrock_client = boto3.client('bedrock-runtime', region_name=self.settings.bedrock_region)
        except Exception as e:
            self.logger.error(f"Failed to initialize Bedrock client: {str(e)}")
            st.error(f"Failed to initialize Bedrock client: {str(e)}")
            self.bedrock_client = None
   
    def invoke_model(self, prompt: str, max_tokens: int = 4000) -> Optional[str]:
        if not self.bedrock_client:
            st.warning("Bedrock client unavailable. Using fallback response.")
            return self._fallback_response(prompt)
       
        try:
            body = {
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": max_tokens,
                "messages": [{"role": "user", "content": prompt}]
            }
            response = self.bedrock_client.invoke_model(
                modelId=self.settings.bedrock_model_id,
                body=json.dumps(body)
            )
            response_body = json.loads(response['body'].read())
            return response_body.get('content', [{}])[0].get('text', '')
        except ClientError as e:
            self.logger.error(f"Bedrock client error: {str(e)}")
            st.error(f"AWS Bedrock Error: {str(e)}")
            return self._fallback_response(prompt)
        except Exception as e:
            self.logger.error(f"Error invoking Bedrock model: {str(e)}")
            st.error(f"Error invoking Bedrock: {str(e)}")
            return self._fallback_response(prompt)
   
    def _fallback_response(self, prompt: str) -> str:
        if "translate" in prompt.lower():
            return "Translation not available - using original text"
        elif "json" in prompt.lower():
            return json.dumps({
                "title": "Generated Presentation",
                "slides": [
                    {
                        "title": "Overview",
                        "content": ["Key insights from the document", "Main topics covered", "Summary of findings"],
                        "type": "content",
                        "chart_suggestion": "none",
                        "chart_data": {}
                    }
                ]
            })
        return "Content analysis not available - using default structure"
   
    def translate_text(self, text: str, target_language: str, source_language: str = 'auto') -> str:
        if source_language == target_language or not text:
            return text
        prompt = f"""
        Translate the following text to {target_language}, maintaining structure and formatting:
        {text[:4000]}
        Translation:
        """
        response = self.invoke_model(prompt)
        return response or text
   
    def generate_presentation_content(self, text: str, company_name: str,
                                     report_month: str, report_year: int,
                                     domain_focus: List[str], language: str,
                                     content_type: str, tables: List[Dict[str, List]],
                                     content_analysis: Dict[str, Any] = None) -> Dict[str, Any]:
        focus_areas = ", ".join(domain_focus)
        table_summary = ""
        if tables:
            table_summary = "\nTables found:\n"
            for i, table in enumerate(tables, 1):
                table_summary += f"Table {i}:\nHeaders: {', '.join(table['headers'])}\nRows: {len(table['rows'])}\n"
       
        prompt = f"""
        Analyze this {content_type} document and create a professional presentation structure for {company_name}
        for {report_month} {report_year} in {language}. Focus areas: {focus_areas}.
       
        Content: {text[:4000]}
        {table_summary}
       
        Steps:
        1. Identify key topics, metrics, and insights from the text and tables.
        2. Group information into 5-7 logical slides, including a title slide.
        3. For each slide, provide a title, 3-5 bullet points, and suggest a chart type (bar, line, pie, table, or none) with sample data.
        4. Tailor content to the {content_type} context, using tables for data-driven slides if available.
       
        Return JSON format:
        {{
            "title": "Presentation title",
            "slides": [
                {{
                    "title": "Slide title",
                    "content": ["point 1", "point 2", "point 3"],
                    "type": "content",
                    "chart_suggestion": "bar|line|pie|table|none",
                    "chart_data": {{"labels": ["A", "B"], "values": [10, 20]}} or {{"headers": [], "rows": [[]]}}
                }}
            ]
        }}
        """
        # Add statistical data to prompt if available
        if content_analysis and content_analysis.get('has_statistics', False):
            statistical_data = content_analysis['statistical_data']
            stats_summary = f"""
            Statistical data found:
            - Numbers: {len(statistical_data['numbers'])} metrics
            - Percentages: {len(statistical_data['percentages'])} values
            - Monetary values: {len(statistical_data['monetary_values'])} amounts
            - Comparisons: {len(statistical_data['comparisons'])} changes
            """
            prompt += stats_summary
       
        response = self.invoke_model(prompt, max_tokens=6000)
        if response:
            try:
                json_start = response.find('{')
                json_end = response.rfind('}') + 1
                json_str = response[json_start:json_end]
                return json.loads(json_str)
            except Exception as e:
                self.logger.error(f"Error parsing JSON: {str(e)}")
                st.error(f"Error parsing Bedrock response: {str(e)}")
        return self._create_default_structure(company_name, report_month, report_year, content_type, tables)
   
    def _create_default_structure(self, company_name: str, report_month: str,
                                 report_year: int, content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
        base_structure = {
            "title": f"{company_name} - {report_month} {report_year} Report",
            "slides": [
                {
                    "title": "Executive Summary",
                    "content": ["Overview of key findings", "Main topics identified", "Summary of performance"],
                    "type": "content",
                    "chart_suggestion": "none",
                    "chart_data": {}
                },
                {
                    "title": "Key Insights",
                    "content": ["Primary observations", "Notable trends", "Actionable takeaways"],
                    "type": "content",
                    "chart_suggestion": "none",
                    "chart_data": {}
                }
            ]
        }
        if content_type == 'logistics':
            base_structure["slides"].append({
                "title": "Logistics Metrics",
                "content": ["Operational efficiency", "Fleet performance", "Delivery metrics"],
                "type": "content",
                "chart_suggestion": "bar",
                "chart_data": {"labels": ["Efficiency", "On-Time", "Cost"], "values": [85, 92, 75]}
            })
        elif content_type == 'business':
            base_structure["slides"].append({
                "title": "Financial Overview",
                "content": ["Revenue trends", "Profit margins", "Cost analysis"],
                "type": "content",
                "chart_suggestion": "line",
                "chart_data": {"labels": ["Q1", "Q2", "Q3"], "values": [100, 120, 130]}
            })
        elif content_type == 'story':
            base_structure["slides"].append({
                "title": "Key Themes",
                "content": ["Main narrative themes", "Character development", "Plot highlights"],
                "type": "content",
                "chart_suggestion": "none",
                "chart_data": {}
            })
        elif content_type == 'sop':
            base_structure["slides"].append({
                "title": "Procedure Overview",
                "content": ["Key procedures", "Compliance requirements", "Implementation steps"],
                "type": "content",
                "chart_suggestion": "table",
                "chart_data": {"headers": ["Step", "Description"], "rows": [["1", "Initiate process"], ["2", "Verify compliance"]]}
            })
        if tables:
            base_structure["slides"].append({
                "title": "Data Summary",
                "content": [f"Table with {len(tables[0]['headers'])} columns and {len(tables[0]['rows'])} rows"],
                "type": "content",
                "chart_suggestion": "table",
                "chart_data": {"headers": tables[0]["headers"], "rows": tables[0]["rows"][:3]}
            })
        return base_structure

class ChartGenerator:
    def __init__(self):
        self.colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
        # Define base layouts for different styles
        self.base_layout_configs = {
            "default": {
                "font_size": 12,
                "title_font_size": 16,
                "bgcolor": 'white',
                "plot_bgcolor": 'white',
                "paper_bgcolor": 'white',
                "margin": dict(l=50, r=50, b=50, t=50),
                "legend_orientation": "h",
                "legend_xanchor": "center",
                "legend_x": 0.5,
                "legend_y": -0.2
            },
            "Minimalist": {
                "font_size": 10,
                "title_font_size": 14,
                "bgcolor": 'rgba(0,0,0,0)',
                "plot_bgcolor": 'rgba(0,0,0,0)',
                "paper_bgcolor": 'rgba(0,0,0,0)',
                "margin": dict(l=30, r=30, b=30, t=30),
                "showlegend": False,
                "xaxis_showgrid": False,
                "yaxis_showgrid": False
            },
            "Creative": {
                "font_size": 14,
                "title_font_size": 20,
                "bgcolor": '#f8f8f8',
                "plot_bgcolor": '#f8f8f8',
                "paper_bgcolor": '#f8f8f8',
                "margin": dict(l=60, r=60, b=60, t=60),
                "hovermode": "x unified"
            }
        }
   
    def create_chart(self, chart_type: str, data: Dict[str, Any], style: str) -> str:
        """Create chart using Plotly and return base64 encoded image"""
        try:
            fig = None
            labels = data.get('labels', ['Category A', 'Category B', 'Category C'])
            values = data.get('values', [10, 20, 30])
            headers = data.get('headers', [])
            rows = data.get('rows', [])
           
            # Apply base style configuration
            layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])

            plt.style.use('seaborn-v0_8')
            if chart_type == 'bar':
                fig = px.bar(x=labels, y=values, color=labels, color_discrete_sequence=self.colors)
                fig.update_layout(title="<b>Data Visualization: Key Metrics</b>", **layout_config)
            elif chart_type == 'line':
                fig = px.line(x=labels, y=values, markers=True, line_shape='spline', color_discrete_sequence=[self.colors[0]])
                fig.update_layout(title="<b>Trend Analysis Over Time</b>", **layout_config)
            elif chart_type == 'pie':
                fig = px.pie(values=values, names=labels, color_discrete_sequence=self.colors, hole=0.3)
                fig.update_layout(title="<b>Distribution Overview</b>", **layout_config)
                fig.update_traces(textinfo='percent+label', pull=[0.05 if i == values.index(max(values)) else 0 for i in range(len(values))])
            elif chart_type == 'table':
                if headers and rows:
                    fig = go.Figure(data=[go.Table(
                        header=dict(values=headers,
                                    fill_color=self.colors[0],
                                    font=dict(color='white', size=layout_config["font_size"] + 2),
                                    align='center',
                                    height=30),
                        cells=dict(values=[rows[i] for i in range(len(rows))],
                                   fill_color='lavender',
                                   font=dict(size=layout_config["font_size"]),
                                   align='left',
                                   height=25)
                    )])
                    fig.update_layout(title="<b>Detailed Data Table</b>", **layout_config)
           
            if fig:
                buffer = io.BytesIO()
                fig.write_image(buffer, format='png', scale=2)
                buffer.seek(0)
                chart_base64 = base64.b64encode(buffer.read()).decode()
                return chart_base64
            return ""
        except Exception as e:
            logger.error(f"Error creating chart: {str(e)}")
            return ""
   
    def create_statistical_chart(self, statistical_data: Dict[str, Any], chart_type: str, style: str) -> str:
        """Create charts from extracted statistical data"""
        try:
            layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])
           
            if chart_type == 'statistics_bar' and statistical_data['numbers']:
                # Create bar chart from numbers
                data = statistical_data['numbers'][:6]  # Limit to 6 items
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
               
                fig = px.bar(x=labels, y=values, color=labels,
                            color_discrete_sequence=self.colors,
                            title="<b>Key Metrics Overview</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
               
            elif chart_type == 'percentage_pie' and statistical_data['percentages']:
                # Create pie chart from percentages
                data = statistical_data['percentages'][:5]  # Limit to 5 items
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
               
                fig = px.pie(values=values, names=labels,
                            color_discrete_sequence=self.colors,
                            title="<b>Percentage Distribution</b>")
                fig.update_layout(**layout_config)
               
            elif chart_type == 'monetary_line' and statistical_data['monetary_values']:
                # Create line chart from monetary values
                data = statistical_data['monetary_values'][:6]
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
               
                fig = px.line(x=labels, y=values, markers=True,
                            color_discrete_sequence=[self.colors[0]],
                            title="<b>Financial Metrics Trend</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
               
            elif chart_type == 'comparison_bar' and statistical_data['comparisons']:
                # Create comparison chart
                data = statistical_data['comparisons'][:6]
                labels = [item['context'] for item in data]
                values = [item['value'] if item['type'] == 'increase' else -item['value'] for item in data]
                colors = ['green' if v > 0 else 'red' for v in values]
               
                fig = px.bar(x=labels, y=values, color=colors,
                            title="<b>Performance Changes (%)</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
               
            else:
                return ""
               
            # Convert to base64
            buffer = io.BytesIO()
            fig.write_image(buffer, format='png', scale=2)
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode()
           
        except Exception as e:
            logger.error(f"Error creating statistical chart: {str(e)}")
            return ""

    def suggest_chart_type(self, statistical_data: Dict[str, Any]) -> str:
        """Suggest the best chart type based on available data"""
        if statistical_data['percentages']:
            return 'percentage_pie'
        elif statistical_data['monetary_values']:
            return 'monetary_line'
        elif statistical_data['comparisons']:
            return 'comparison_bar'
        elif statistical_data['numbers']:
            return 'statistics_bar'
        return 'none'

class PresentationGenerator:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.logger = logging.getLogger(__name__)
        self.chart_generator = ChartGenerator()
        self.default_theme = {
            "colors": {
                "primary": "#003087",  # Dark blue
                "secondary": "#FFFFFF",  # White
                "accent": "#FFA500",  # Orange
                "text": "#404040",  # Dark gray
                "footer": "#808080"  # Gray
            },
            "fonts": {
                "title": {"name": "Calibri", "size": 44, "bold": True},
                "subtitle": {"name": "Calibri", "size": 24, "bold": False},
                "content": {"name": "Calibri", "size": 18, "bold": False},
                "footer": {"name": "Calibri", "size": 10, "bold": False}
            },
            "footer": {
                "text": "{company_name} | {date} | Confidential",
                "position": {"left": 0.5, "top": 7.0, "width": 9.0, "height": 0.5}
            }
        }
        self.theme = self.default_theme
        try:
            theme_path = Path(__file__).parent / "templates" / "styles" / "corporate_theme.json"
            if theme_path.exists():
                with open(theme_path, "r") as f:
                    loaded_theme = json.load(f)
                    self.theme = {**self.default_theme, **loaded_theme.get("theme", {})}
                    for key in ["colors", "fonts", "footer"]:
                        if key in loaded_theme.get("theme", {}) and isinstance(loaded_theme["theme"][key], dict):
                            self.theme[key] = {**self.default_theme[key], **loaded_theme["theme"][key]}
            else:
                self.logger.warning(f"Theme file not found at {theme_path}. Using default theme.")
        except Exception as e:
            self.logger.warning(f"Failed to load theme: {str(e)}. Using default theme.")
       
        self.primary_color = RGBColor.from_string(self.theme["colors"]["primary"].lstrip("#"))
        self.secondary_color = RGBColor.from_string(self.theme["colors"]["secondary"].lstrip("#"))
        self.accent_color = RGBColor.from_string(self.theme["colors"]["accent"].lstrip("#"))
        self.text_color = RGBColor.from_string(self.theme["colors"]["text"].lstrip("#"))
        self.footer_color = RGBColor.from_string(self.theme["colors"]["footer"].lstrip("#"))
   
    def generate_presentation(self, content: str, company_name: str,
                             report_month: str, report_year: int,
                             domain_focus: List[str], language: str,
                             content_type: str, tables: List[Dict[str, List]],
                             style: str, content_analysis: Dict[str, Any] = None) -> str:
        try:
            # Initialize presentation based on template choice
            if self.settings.use_custom_template and self.settings.template_path and self.settings.template_path.exists():
                try:
                    prs = Presentation(str(self.settings.template_path))
                    self.logger.info(f"Successfully loaded custom template from {self.settings.template_path}")
                except Exception as e:
                    self.logger.warning(f"Failed to load custom template: {str(e)}. Using blank presentation.")
                    prs = Presentation()
            else:
                prs = Presentation()
            
            # Set slide size to widescreen (16:9) if not using template
            if not self.settings.use_custom_template:
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                # Set master slide background
                background = prs.slide_master.background
                background.fill.solid()
                background.fill.fore_color.rgb = self.secondary_color

            bedrock_client = BedrockClient(self.settings)
            content_structure = bedrock_client.generate_presentation_content(
                content, company_name, report_month, report_year, domain_focus, language, content_type, tables, content_analysis
            )
            if not content_structure:
                content_structure = bedrock_client._create_default_structure(company_name, report_month, report_year, content_type, tables)
           
            self._create_enhanced_title_slide(prs, content_structure['title'], company_name, report_month, report_year, domain_focus, style)
            for slide_data in content_structure['slides']:
                self._create_enhanced_content_slide(prs, slide_data, style)
           
            # Add statistical slide if applicable
            if content_analysis and content_analysis.get('has_statistics', False):
                self._create_statistical_slide(prs, content_analysis['statistical_data'], style)
           
            filename = f"Presentation_{company_name.replace(' ', '_')}_{report_month}_{report_year}.pptx"
            output_path = self.settings.output_folder / filename
            prs.save(str(output_path))
            return str(output_path)
        except Exception as e:
            self.logger.error(f"Error generating presentation: {str(e)}")
            st.error(f"Error generating presentation: {str(e)}")
            raise
    
    def _create_enhanced_title_slide(self, prs: Presentation, title: str, company_name: str,
                                    report_month: str, report_year: int, domain_focus: List[str], style: str):
        # Create blank slide (no layout)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout
       
        # Add background rectangle
        left, top, width, height = Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        background.fill.solid()
        background.fill.fore_color.rgb = self.primary_color
        background.line.fill.background()
       
        # Add title
        title_left, title_top, title_width, title_height = Inches(1), Inches(2), Inches(11.333), Inches(2)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
        title_paragraph.font.color.rgb = self.secondary_color
        title_paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
        title_paragraph.alignment = PP_ALIGN.CENTER
       
        # Add subtitle
        subtitle_text = f"{company_name} | {report_month} {report_year} | {', '.join(domain_focus)}"
        subtitle_top = Inches(4.5)
        subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
        subtitle_paragraph.font.color.rgb = self.accent_color
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
       
        # Add decorative elements based on style
        if style == "Creative":
            self._add_creative_decorations(slide)
        elif style == "Minimalist":
            self._add_minimalist_decorations(slide)
       
        self._add_footer(slide, company_name, f"{report_month} {report_year}")
   
    def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], style: str):
        # Create blank slide (no layout)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout
       
        # Add title
        title_left, title_top, title_width, title_height = Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"] - 10)  # Smaller than main title
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
       
        # Add content divider
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(0.1))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.accent_color
        divider.line.fill.background()
       
        # Populate content (bullet points or chart)
        if slide_data.get('chart_suggestion') != "none" and slide_data.get('chart_data'):
            # Content on left, chart on right
            content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(5.5), Inches(5)
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            self._add_bullet_points(content_box.text_frame, slide_data['content'])

            chart_left = Inches(6.5)
            self._add_chart_to_slide(slide, slide_data['chart_suggestion'], slide_data['chart_data'], chart_left, content_top, style)
        else:
            # Full width content
            content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(12), Inches(5)
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            self._add_bullet_points(content_box.text_frame, slide_data['content'])
       
        self._add_footer(slide, "Your Company", f"{datetime.now().strftime('%B %Y')}")
   
    def _create_statistical_slide(self, prs: Presentation, statistical_data: Dict[str, Any], style: str):
        """Create a slide specifically for statistical visualizations"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout
       
        # Add title
        title_left, title_top, title_width, title_height = Inches(0.5), Inches(0.5), Inches(12), Inches(1)
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        title_frame.text = "📊 Statistical Analysis"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"] - 10)
        title_paragraph.font.color.rgb = self.primary_color
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
       
        # Add divider
        divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12), Inches(0.1))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.accent_color
        divider.line.fill.background()
       
        # Add statistical charts
        chart_type = self.chart_generator.suggest_chart_type(statistical_data)
        if chart_type != 'none':
            chart_base64 = self.chart_generator.create_statistical_chart(statistical_data, chart_type, style)
            if chart_base64:
                buffer = io.BytesIO(base64.b64decode(chart_base64))
                slide.shapes.add_picture(buffer, Inches(1), Inches(1.5), Inches(8), Inches(4.5))
       
        # Add statistical summary text
        summary_text = self._generate_statistical_summary(statistical_data)
        if summary_text:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = summary_text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = self.text_color
       
        self._add_footer(slide, "Statistical Analysis", f"{datetime.now().strftime('%B %Y')}")
   
    def _generate_statistical_summary(self, statistical_data: Dict[str, Any]) -> str:
        """Generate a text summary of statistical findings"""
        summary_parts = []
       
        if statistical_data['numbers']:
            summary_parts.append(f"Key metrics: {len(statistical_data['numbers'])} data points identified")
       
        if statistical_data['percentages']:
            avg_percentage = sum(item['value'] for item in statistical_data['percentages']) / len(statistical_data['percentages'])
            summary_parts.append(f"Average percentage value: {avg_percentage:.1f}%")
       
        if statistical_data['monetary_values']:
            total_value = sum(item['value'] for item in statistical_data['monetary_values'])
            summary_parts.append(f"Total monetary value: ${total_value:,.0f}")
       
        return " | ".join(summary_parts)
   
    def _add_bullet_points(self, text_frame, content_list: List[str]):
        text_frame.clear()
        for i, item in enumerate(content_list):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(self.theme["fonts"]["content"]["size"])
            paragraph.font.color.rgb = self.text_color
            paragraph.level = 0
            paragraph.space_after = Pt(12)
   
    def _add_chart_to_slide(self, slide, chart_type: str, chart_data: Dict[str, Any], left: float, top: float, style: str):
        try:
            chart_width, chart_height = Inches(6), Inches(4)
            if chart_type == 'table':
                table = slide.shapes.add_table(
                    len(chart_data.get('rows', [])[:5]) + 1,
                    len(chart_data.get('headers', [])),
                    left, top, chart_width, chart_height
                ).table
                for i, header in enumerate(chart_data.get('headers', [])):
                    table.cell(0, i).text = header
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)
                    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = self.secondary_color
                    table.cell(0, i).fill.solid()
                    table.cell(0, i).fill.fore_color.rgb = self.primary_color
                for i, row in enumerate(chart_data.get('rows', [])[:5], 1):
                    for j, value in enumerate(row):
                        table.cell(i, j).text = str(value)
                        table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
                        table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = self.text_color
                        if i % 2 == 0:
                            table.cell(i, j).fill.solid()
                            table.cell(i, j).fill.fore_color.rgb = RGBColor(240, 240, 240)
            else:
                chart_base64 = self.chart_generator.create_chart(chart_type, chart_data, style)
                if chart_base64:
                    buffer = io.BytesIO(base64.b64decode(chart_base64))
                    slide.shapes.add_picture(buffer, left, top, chart_width, chart_height)
        except Exception as e:
            self.logger.error(f"Error adding chart: {str(e)}")
            st.error(f"Error adding chart: {str(e)}")
            fallback_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
            fallback_box.text_frame.text = "📊 Chart unavailable\n(Data visualization)"
            fallback_box.text_frame.paragraphs[0].font.size = Pt(16)
            fallback_box.text_frame.paragraphs[0].font.color.rgb = self.accent_color
   
    def _add_creative_decorations(self, slide):
        """Add creative decorative elements to slide"""
        try:
            # Add accent shapes
            for i in range(3):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.5 + i*2), Inches(6), Inches(1.5), Inches(0.3)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.accent_color
                shape.line.fill.background()
               
            # Add small decorative icons
            icon_path = Path(__file__).parent / "templates" / "icons" / "creative_icon.png"
            if not icon_path.exists():
                # Create a simple colored circle if icon doesn't exist
                img = Image.new('RGB', (200, 200), color=(self.accent_color.rgb[0], self.accent_color.rgb[1], self.accent_color.rgb[2]))
                img.save(icon_path)
           
            slide.shapes.add_picture(str(icon_path), Inches(11), Inches(0.5), Inches(1.5), Inches(1.5))
           
        except Exception as e:
            self.logger.error(f"Error adding creative decorations: {str(e)}")

    def _add_minimalist_decorations(self, slide):
        """Add minimalist decorative elements to slide"""
        try:
            # Add thin line at bottom
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(7), Inches(13.333), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.primary_color
            line.line.fill.background()
        except Exception as e:
            self.logger.error(f"Error adding minimalist decorations: {str(e)}")

    def _add_footer(self, slide, company_name: str, date: str):
        left = Inches(self.theme["footer"]["position"]["left"])
        top = Inches(self.theme["footer"]["position"]["top"])
        width = Inches(self.theme["footer"]["position"]["width"])
        height = Inches(self.theme["footer"]["position"]["height"])

        footer_box = slide.shapes.add_textbox(left, top, width, height)
        footer_frame = footer_box.text_frame
        footer_frame.text = self.theme["footer"]["text"].format(company_name=company_name, date=date)
        footer_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
        footer_frame.paragraphs[0].font.color.rgb = self.footer_color
        footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add slide number
        slide_number_left = Inches(0.5)
        slide_number_box = slide.shapes.add_textbox(slide_number_left, top, Inches(1), height)
        slide_number_frame = slide_number_box.text_frame
        slide_number_frame.text = f"Page {slide.slide_id}"
        slide_number_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
        slide_number_frame.paragraphs[0].font.color.rgb = self.footer_color
        slide_number_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

@st.cache_resource
def initialize_components():
    settings = Settings()
    doc_processor = DocumentProcessor()
    presentation_generator = PresentationGenerator(settings)
    return settings, doc_processor, presentation_generator

def main():
    st.markdown("""
    <div class="main-header">
        <h1>🎯 Smart Presentation Generator</h1>
        <p>Transform any content into beautiful, professional presentations with AI-powered insights and visualizations</p>
    </div>
    """, unsafe_allow_html=True)
   
    settings, doc_processor, presentation_generator = initialize_components()
    content_type = "Not yet processed"
   
    with st.sidebar:
        st.header("⚙️ Configuration")
        input_method = st.radio(
            "Choose input method:",
            ["📄 Upload PDF", "📝 Direct Text Input"],
            horizontal=True
        )
        
        # Template selection section
        st.markdown("---")
        st.subheader("🎨 Template Options")
        template_option = st.radio(
            "Select template source:",
            ["🚀 Use Default Template", "📁 Upload Custom Template", "🌐 Enter Template URL"],
            index=0
        )
        
        custom_template = None
        if template_option == "📁 Upload Custom Template":
            custom_template = st.file_uploader("Upload PowerPoint Template", type=["pptx"])
            if custom_template:
                # Save the uploaded template
                template_path = settings.template_folder / custom_template.name
                with open(template_path, "wb") as f:
                    f.write(custom_template.getvalue())
                settings.set_template_path(template_path)
        elif template_option == "🌐 Enter Template URL":
            template_url = st.text_input("Enter Template URL (must be direct .pptx file link)")
            if template_url:
                settings.set_template_path(template_url)
        else:
            settings.set_template_path(None)  # Use default
            
        st.markdown("---")
        target_language = st.selectbox(
            "Target Language",
            options=["en", "es", "fr", "de", "it", "pt", "nl", "ru", "zh", "ja", "ko", "ar", "hi"],
            format_func=lambda x: {
                "en": "🇺🇸 English", "es": "🇪🇸 Spanish", "fr": "🇫🇷 French",
                "de": "🇩🇪 German", "it": "🇮🇹 Italian", "pt": "🇵🇹 Portuguese",
                "nl": "🇳🇱 Dutch", "ru": "🇷🇺 Russian", "zh": "🇨🇳 Chinese",
                "ja": "🇯🇵 Japanese", "ko": "🇰🇷 Korean", "ar": "🇸🇦 Arabic",
                "hi": "🇮🇳 Hindi"
            }[x]
        )
        company_name = st.text_input("Company Name", "Your Company")
        col1, col2 = st.columns(2)
        with col1:
            report_month = st.selectbox(
                "Month",
                options=["January", "February", "March", "April", "May", "June",
                         "July", "August", "September", "October", "November", "December"],
                index=datetime.now().month - 1
            )
        with col2:
            report_year = st.number_input("Year", min_value=2020, max_value=2030, value=datetime.now().year)
        domain_focus = st.multiselect(
            "Focus Areas",
            options=["Logistics Vehicles", "Cargo Management", "Fleet Operations",
                     "Supply Chain", "Transportation Efficiency", "Business Performance",
                     "Financial Reporting", "Market Analysis", "Strategic Planning",
                     "Narrative Development", "Character Analysis", "Plot Structure",
                     "Standard Operating Procedures", "Compliance and Regulations", "Operational Guidelines"],
            default=["Business Performance"]
        )
        presentation_style = st.selectbox(
            "Presentation Style",
            options=["Default", "Minimalist", "Creative"]
        )
       
        st.markdown("---")
        st.markdown("Developed by **Your Team Name**")

    document_content = ""
    tables_extracted: List[Dict[str, List]] = []
    content_analysis = None

    if input_method == "📄 Upload PDF":
        uploaded_file = st.file_uploader("Upload a PDF document", type=["pdf"])
        if uploaded_file:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                temp_pdf_path = tmp_file.name
           
            with st.spinner("Extracting text and analyzing content..."):
                document_content = doc_processor.extract_text_from_pdf(temp_pdf_path)
                tables_extracted = doc_processor.extract_tables(document_content)
                content_analysis = doc_processor.analyze_content_type(document_content)
                content_type = content_analysis['type']
                st.session_state['document_content'] = document_content
                st.session_state['tables_extracted'] = tables_extracted
                st.session_state['content_type'] = content_type
                st.session_state['content_analysis'] = content_analysis
           
            os.unlink(temp_pdf_path)
            if document_content:
                st.success("PDF processed successfully!")
                st.expander("Extracted Text Preview").write(document_content[:1000] + "...")
                if tables_extracted:
                    st.expander("Extracted Tables Preview").write(tables_extracted)
            else:
                st.error("Failed to extract content from PDF.")
   
    elif input_method == "📝 Direct Text Input":
        text_input = st.text_area("Paste your content here:", height=300)
        if text_input:
            with st.spinner("Analyzing content..."):
                document_content = doc_processor.process_text_document(text_input)
                tables_extracted = doc_processor.extract_tables(document_content)
                content_analysis = doc_processor.analyze_content_type(document_content)
                content_type = content_analysis['type']
                st.session_state['document_content'] = document_content
                st.session_state['tables_extracted'] = tables_extracted
                st.session_state['content_type'] = content_type
                st.session_state['content_analysis'] = content_analysis
            st.success("Text content processed successfully!")
   
    # Display statistical analysis in UI
    if 'content_analysis' in st.session_state and st.session_state['content_analysis'].get('has_statistics', False):
        st.subheader("📊 Statistical Analysis Found")
       
        statistical_data = st.session_state['content_analysis']['statistical_data']
       
        col1, col2 = st.columns(2)
       
        with col1:
            if statistical_data['numbers']:
                st.write("**Key Metrics:**")
                for item in statistical_data['numbers'][:5]:
                    st.write(f"• {item['context']}: {item['original']}")
       
        with col2:
            if statistical_data['percentages']:
                st.write("**Percentages:**")
                for item in statistical_data['percentages'][:5]:
                    st.write(f"• {item['context']}: {item['original']}")
        
        # Show chart previews for all possible chart types
        chart_generator = ChartGenerator()
        st.subheader("📈 Potential Visualizations")
        
        if statistical_data['numbers']:
            st.write("**Bar Chart Preview**")
            chart_base64 = chart_generator.create_statistical_chart(statistical_data, 'statistics_bar', "Default")
            if chart_base64:
                st.image(base64.b64decode(chart_base64))
        
        if statistical_data['percentages']:
            st.write("**Pie Chart Preview**")
            chart_base64 = chart_generator.create_statistical_chart(statistical_data, 'percentage_pie', "Default")
            if chart_base64:
                st.image(base64.b64decode(chart_base64))
        
        if statistical_data['monetary_values']:
            st.write("**Line Chart Preview**")
            chart_base64 = chart_generator.create_statistical_chart(statistical_data, 'monetary_line', "Default")
            if chart_base64:
                st.image(base64.b64decode(chart_base64))
        
        if statistical_data['comparisons']:
            st.write("**Comparison Chart Preview**")
            chart_base64 = chart_generator.create_statistical_chart(statistical_data, 'comparison_bar', "Default")
            if chart_base64:
                st.image(base64.b64decode(chart_base64))

    if st.button("Generate Presentation 🚀", type="primary"):
        if 'document_content' not in st.session_state or not st.session_state['document_content']:
            st.error("Please upload a document or enter text content first.")
            return

        with st.spinner("Generating professional presentation... This may take a few moments."):
            try:
                # Retrieve processed content from session state
                doc_content = st.session_state['document_content']
                tables = st.session_state['tables_extracted']
                c_type = st.session_state['content_type']
                content_analysis = st.session_state.get('content_analysis', {})

                # Translate content if target language is not English
                if target_language != 'en':
                    bedrock_client = BedrockClient(settings)
                    doc_content = bedrock_client.translate_text(doc_content, target_language)
               
                output_ppt_path = presentation_generator.generate_presentation(
                    doc_content, company_name, report_month, report_year, domain_focus, target_language,
                    c_type, tables, presentation_style, content_analysis
                )
               
                if output_ppt_path and os.path.exists(output_ppt_path):
                    with open(output_ppt_path, "rb") as file:
                        btn = st.download_button(
                            label="Download Presentation ⬇️",
                            data=file,
                            file_name=os.path.basename(output_ppt_path),
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    st.success("Presentation generated successfully! Click the button above to download.")
                    st.balloons()
                else:
                    st.error("Failed to generate presentation. Please check logs for details.")
            except Exception as e:
                st.error(f"An unexpected error occurred during presentation generation: {e}")
                logger.error(f"Critical error in main generation flow: {e}")

    st.markdown("---")
    st.info("💡 **Tips for best results:**\n"
            "- Upload clear, well-formatted PDF documents.\n"
            "- Specify relevant focus areas for tailored content.\n"
            "- For custom templates, use a clean template with minimal placeholders.\n"
            "- Ensure your AWS Bedrock credentials are set up in `.env` if you're using LLM functionality.")

if __name__ == "__main__":
    main()