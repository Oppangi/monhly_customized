
##############code 1
# import streamlit as st
# import os
# import tempfile
# import json
# import logging
# from pathlib import Path
# from datetime import datetime
# from typing import Dict, Any, List, Optional, Union
# import zipfile
# import io
# import base64

# # Core libraries
# import boto3
# from botocore.exceptions import ClientError
# import PyPDF2
# from langdetect import detect
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.enum.shapes import MSO_SHAPE
# import pandas as pd
# import matplotlib.pyplot as plt
# import seaborn as sns
# from dotenv import load_dotenv
# import plotly.graph_objects as go
# import plotly.express as px
# from plotly.subplots import make_subplots
# import numpy as np

# # Configure page
# st.set_page_config(
#     page_title="🎯 Smart Presentation Generator",
#     page_icon="🎯",
#     layout="wide",
#     initial_sidebar_state="expanded"
# )

# # Custom CSS for better styling
# st.markdown("""
# <style>
#     .main-header {
#         background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
#         padding: 20px;
#         border-radius: 10px;
#         color: white;
#         margin-bottom: 30px;
#         text-align: center;
#     }
#     .feature-box {
#         background: #f8f9fa;
#         padding: 15px;
#         border-radius: 8px;
#         border-left: 4px solid #667eea;
#         margin: 10px 0;
#     }
#     .metric-container {
#         background: white;
#         padding: 20px;
#         border-radius: 10px;
#         box-shadow: 0 2px 4px rgba(0,0,0,0.1);
#         margin: 10px 0;
#     }
# </style>
# """, unsafe_allow_html=True)

# class Settings:
#     def __init__(self):
#         load_dotenv()
#         self.aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID')
#         self.aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY')
#         self.aws_session_token = os.getenv('AWS_SESSION_TOKEN')
#         self.aws_default_region = os.getenv('AWS_DEFAULT_REGION', 'us-east-1')
#         self.bedrock_model_id = os.getenv('BEDROCK_MODEL_ID', 'anthropic.claude-3-sonnet-20240229-v1:0')
#         self.bedrock_region = os.getenv('BEDROCK_REGION', 'us-east-1')
#         self.upload_folder = Path('uploads')
#         self.output_folder = Path('output')
#         self.max_file_size = 10485760  # 10MB
        
#         # Create directories
#         self.upload_folder.mkdir(exist_ok=True)
#         self.output_folder.mkdir(exist_ok=True)
        
#         # Set up AWS session
#         try:
#             boto3.setup_default_session(
#                 aws_access_key_id=self.aws_access_key_id,
#                 aws_secret_access_key=self.aws_secret_access_key,
#                 aws_session_token=self.aws_session_token,
#                 region_name=self.aws_default_region
#             )
#         except Exception as e:
#             st.warning(f"AWS setup warning: {str(e)}")

# class DocumentProcessor:
#     def __init__(self):
#         self.logger = logging.getLogger(__name__)
    
#     def extract_text_from_pdf(self, pdf_path: str) -> str:
#         """Extract text from PDF file"""
#         try:
#             text = ""
#             with open(pdf_path, 'rb') as file:
#                 pdf_reader = PyPDF2.PdfReader(file)
#                 for page in pdf_reader.pages:
#                     text += page.extract_text()
#             return self._clean_text(text)
#         except Exception as e:
#             self.logger.error(f"Error extracting text from PDF: {str(e)}")
#             return ""
    
#     def process_text_document(self, content: str) -> str:
#         """Process raw text content"""
#         return self._clean_text(content)
    
#     def _clean_text(self, text: str) -> str:
#         """Clean and normalize text"""
#         text = ' '.join(text.split())
#         text = text.replace('\x00', '').replace('\ufeff', '')
#         return text
    
#     def detect_language(self, text: str) -> Optional[str]:
#         """Detect language of text"""
#         try:
#             if len(text.strip()) < 10:
#                 return None
#             return detect(text)
#         except:
#             return None
    
#     def analyze_content_type(self, text: str) -> Dict[str, Any]:
#         """Analyze content to determine type and extract insights"""
#         content_analysis = {
#             'type': 'general',
#             'has_numbers': bool(any(char.isdigit() for char in text)),
#             'has_dates': bool(any(word for word in text.split() if any(char.isdigit() for char in word) and len(word) > 3)),
#             'topics': [],
#             'potential_charts': []
#         }
        
#         # Detect content type
#         logistics_keywords = ['logistics', 'cargo', 'fleet', 'vehicle', 'transport', 'shipping', 'warehouse']
#         business_keywords = ['revenue', 'profit', 'sales', 'market', 'customer', 'performance']
#         story_keywords = ['story', 'narrative', 'character', 'plot', 'once upon']
        
#         text_lower = text.lower()
        
#         if any(keyword in text_lower for keyword in logistics_keywords):
#             content_analysis['type'] = 'logistics'
#         elif any(keyword in text_lower for keyword in business_keywords):
#             content_analysis['type'] = 'business'
#         elif any(keyword in text_lower for keyword in story_keywords):
#             content_analysis['type'] = 'story'
        
#         # Extract potential chart data
#         if content_analysis['has_numbers']:
#             content_analysis['potential_charts'] = ['bar', 'line', 'pie']
        
#         return content_analysis

# class BedrockClient:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         try:
#             self.bedrock_client = boto3.client(
#                 'bedrock-runtime',
#                 region_name=self.settings.bedrock_region
#             )
#         except Exception as e:
#             self.logger.error(f"Failed to initialize Bedrock client: {str(e)}")
#             self.bedrock_client = None
    
#     def invoke_model(self, prompt: str, max_tokens: int = 4000) -> Optional[str]:
#         """Invoke Bedrock model with fallback"""
#         if not self.bedrock_client:
#             return self._fallback_response(prompt)
        
#         try:
#             body = {
#                 "anthropic_version": "bedrock-2023-05-31",
#                 "max_tokens": max_tokens,
#                 "messages": [{"role": "user", "content": prompt}]
#             }
#             response = self.bedrock_client.invoke_model(
#                 modelId=self.settings.bedrock_model_id,
#                 body=json.dumps(body)
#             )
#             response_body = json.loads(response['body'].read())
#             return response_body.get('content', [{}])[0].get('text', '')
#         except Exception as e:
#             self.logger.error(f"Bedrock error: {str(e)}")
#             return self._fallback_response(prompt)
    
#     def _fallback_response(self, prompt: str) -> str:
#         """Fallback response when Bedrock is unavailable"""
#         if "translate" in prompt.lower():
#             return "Translation not available - using original text"
#         elif "json" in prompt.lower():
#             return json.dumps({
#                 "title": "Generated Presentation",
#                 "slides": [
#                     {
#                         "title": "Overview",
#                         "content": ["Key insights from the document", "Main topics covered", "Summary of findings"],
#                         "type": "content"
#                     }
#                 ]
#             })
#         return "Content analysis not available - using default structure"
    
#     def generate_presentation_content(self, text: str, company_name: str, 
#                                     report_month: str, report_year: int, 
#                                     domain_focus: List[str], language: str,
#                                     content_type: str = 'general') -> Dict[str, Any]:
#         """Generate presentation content with enhanced analysis"""
#         focus_areas = ", ".join(domain_focus)
        
#         if content_type == 'logistics':
#             prompt = self._get_logistics_prompt(text, company_name, report_month, report_year, focus_areas, language)
#         elif content_type == 'business':
#             prompt = self._get_business_prompt(text, company_name, report_month, report_year, language)
#         elif content_type == 'story':
#             prompt = self._get_story_prompt(text, language)
#         else:
#             prompt = self._get_general_prompt(text, company_name, report_month, report_year, language)
        
#         response = self.invoke_model(prompt, max_tokens=6000)
        
#         if response:
#             try:
#                 json_start = response.find('{')
#                 json_end = response.rfind('}') + 1
#                 json_str = response[json_start:json_end]
#                 return json.loads(json_str)
#             except Exception as e:
#                 self.logger.error(f"Error parsing JSON: {str(e)}")
        
#         return self._create_default_structure(company_name, report_month, report_year, content_type)
    
#     def _get_logistics_prompt(self, text: str, company_name: str, report_month: str, 
#                              report_year: int, focus_areas: str, language: str) -> str:
#         return f"""
#         Analyze this logistics document and create a comprehensive presentation structure for {company_name}.
#         Focus areas: {focus_areas}
#         Language: {language}
        
#         Content: {text[:4000]}
        
#         Create a JSON structure with:
#         1. Title slide
#         2. Executive Summary with key metrics
#         3. Logistics Performance (include chart suggestions)
#         4. Fleet Analysis (include chart suggestions)
#         5. Challenges and Solutions
#         6. Future Outlook
        
#         For each slide with data, suggest appropriate charts (bar, line, pie, scatter).
        
#         Return JSON format:
#         {{
#             "title": "Presentation title",
#             "slides": [
#                 {{
#                     "title": "Slide title",
#                     "content": ["point 1", "point 2", "point 3"],
#                     "type": "content",
#                     "chart_suggestion": "bar|line|pie|scatter",
#                     "chart_data": {{"labels": ["A", "B"], "values": [10, 20]}}
#                 }}
#             ]
#         }}
#         """
    
#     def _get_business_prompt(self, text: str, company_name: str, report_month: str, 
#                             report_year: int, language: str) -> str:
#         return f"""
#         Analyze this business document and create a presentation for {company_name}.
#         Language: {language}
        
#         Content: {text[:4000]}
        
#         Create slides covering:
#         1. Executive Summary
#         2. Key Performance Indicators
#         3. Market Analysis
#         4. Financial Overview
#         5. Strategic Initiatives
        
#         Include chart suggestions for data visualization.
        
#         Return JSON format with chart suggestions and sample data.
#         """
    
#     def _get_story_prompt(self, text: str, language: str) -> str:
#         return f"""
#         Analyze this story/narrative content and create an engaging presentation.
#         Language: {language}
        
#         Content: {text[:4000]}
        
#         Create slides covering:
#         1. Story Overview
#         2. Main Characters/Elements
#         3. Key Themes
#         4. Plot Development
#         5. Conclusion/Lessons
        
#         Make it visually engaging with creative slide titles.
        
#         Return JSON format.
#         """
    
#     def _get_general_prompt(self, text: str, company_name: str, report_month: str, 
#                            report_year: int, language: str) -> str:
#         return f"""
#         Analyze this document and create a professional presentation.
#         Language: {language}
        
#         Content: {text[:4000]}
        
#         Create 5-7 slides based on the content's main topics.
#         Include data visualization suggestions where appropriate.
        
#         Return JSON format with chart suggestions.
#         """
    
#     def _create_default_structure(self, company_name: str, report_month: str, 
#                                  report_year: int, content_type: str) -> Dict[str, Any]:
#         """Create default presentation structure based on content type"""
#         if content_type == 'logistics':
#             return {
#                 "title": f"{company_name} - Logistics Report {report_month} {report_year}",
#                 "slides": [
#                     {
#                         "title": "Executive Summary",
#                         "content": ["Monthly performance overview", "Key operational metrics", "Strategic achievements"],
#                         "type": "content",
#                         "chart_suggestion": "bar",
#                         "chart_data": {"labels": ["Performance", "Efficiency", "Quality"], "values": [85, 92, 78]}
#                     },
#                     {
#                         "title": "Fleet Performance",
#                         "content": ["Vehicle utilization rates", "Maintenance efficiency", "Route optimization"],
#                         "type": "content",
#                         "chart_suggestion": "line",
#                         "chart_data": {"labels": ["Jan", "Feb", "Mar"], "values": [80, 85, 90]}
#                     }
#                 ]
#             }
#         elif content_type == 'story':
#             return {
#                 "title": "Story Presentation",
#                 "slides": [
#                     {
#                         "title": "Story Overview",
#                         "content": ["Main narrative elements", "Setting and context", "Key themes"],
#                         "type": "content"
#                     }
#                 ]
#             }
#         else:
#             return {
#                 "title": f"{company_name} - Presentation {report_month} {report_year}",
#                 "slides": [
#                     {
#                         "title": "Overview",
#                         "content": ["Key insights", "Main findings", "Summary points"],
#                         "type": "content"
#                     }
#                 ]
#             }

# class ChartGenerator:
#     def __init__(self):
#         self.colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
    
#     def create_sample_chart(self, chart_type: str, data: Dict[str, Any]) -> str:
#         """Create sample chart and return base64 encoded image"""
#         try:
#             plt.figure(figsize=(10, 6))
#             plt.style.use('seaborn-v0_8')
            
#             labels = data.get('labels', ['A', 'B', 'C'])
#             values = data.get('values', [10, 20, 30])
            
#             if chart_type == 'bar':
#                 plt.bar(labels, values, color=self.colors[:len(labels)])
#                 plt.title('Sample Bar Chart')
#             elif chart_type == 'line':
#                 plt.plot(labels, values, marker='o', linewidth=2, color=self.colors[0])
#                 plt.title('Sample Line Chart')
#             elif chart_type == 'pie':
#                 plt.pie(values, labels=labels, autopct='%1.1f%%', colors=self.colors[:len(labels)])
#                 plt.title('Sample Pie Chart')
            
#             plt.tight_layout()
            
#             # Convert to base64
#             buffer = io.BytesIO()
#             plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
#             buffer.seek(0)
#             chart_base64 = base64.b64encode(buffer.read()).decode()
#             plt.close()
            
#             return chart_base64
#         except Exception as e:
#             logging.error(f"Error creating chart: {str(e)}")
#             return ""

# class PresentationGenerator:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         self.chart_generator = ChartGenerator()
        
#         # Corporate colors
#         self.primary_color = RGBColor(0, 48, 135)  # Dark blue
#         self.secondary_color = RGBColor(255, 255, 255)  # White
#         self.accent_color = RGBColor(102, 126, 234)  # Purple
#         self.text_color = RGBColor(64, 64, 64)  # Dark gray
    
#     def generate_presentation(self, content: str, company_name: str, 
#                             report_month: str, report_year: int, 
#                             domain_focus: List[str], language: str,
#                             content_type: str = 'general') -> str:
#         """Generate enhanced PowerPoint presentation"""
#         try:
#             # Create presentation
#             prs = Presentation()
            
#             # Generate content structure
#             bedrock_client = BedrockClient(self.settings)
#             content_structure = bedrock_client.generate_presentation_content(
#                 content, company_name, report_month, report_year, 
#                 domain_focus, language, content_type
#             )
            
#             # Create slides
#             self._create_enhanced_title_slide(prs, content_structure['title'], 
#                                             company_name, report_month, report_year)
            
#             for slide_data in content_structure['slides']:
#                 self._create_enhanced_content_slide(prs, slide_data)
            
#             # Save presentation
#             filename = f"Enhanced_Presentation_{company_name.replace(' ', '_')}_{report_month}_{report_year}.pptx"
#             output_path = self.settings.output_folder / filename
#             prs.save(str(output_path))
            
#             return str(output_path)
        
#         except Exception as e:
#             self.logger.error(f"Error generating presentation: {str(e)}")
#             raise
    
#     def _create_enhanced_title_slide(self, prs: Presentation, title: str, 
#                                    company_name: str, report_month: str, report_year: int):
#         """Create enhanced title slide with modern design"""
#         slide_layout = prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
        
#         # Remove default placeholders
#         for shape in slide.shapes:
#             if hasattr(shape, 'text'):
#                 shape.text = ""
        
#         # Add custom title
#         left = Inches(1)
#         top = Inches(2)
#         width = Inches(8)
#         height = Inches(2)
        
#         title_box = slide.shapes.add_textbox(left, top, width, height)
#         title_frame = title_box.text_frame
#         title_frame.text = title
        
#         # Style title
#         title_paragraph = title_frame.paragraphs[0]
#         title_paragraph.font.size = Pt(48)
#         title_paragraph.font.color.rgb = self.primary_color
#         title_paragraph.font.bold = True
#         title_paragraph.alignment = PP_ALIGN.CENTER
        
#         # Add subtitle
#         subtitle_top = Inches(4.5)
#         subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1.5))
#         subtitle_frame = subtitle_box.text_frame
#         subtitle_frame.text = f"{company_name} | {report_month} {report_year}"
        
#         subtitle_paragraph = subtitle_frame.paragraphs[0]
#         subtitle_paragraph.font.size = Pt(24)
#         subtitle_paragraph.font.color.rgb = self.accent_color
#         subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
#         # Add decorative shape
#         self._add_decorative_shape(slide)
    
#     def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
#         """Create enhanced content slide with charts if applicable"""
#         slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(slide_layout)
        
#         # Title
#         title_shape = slide.shapes.title
#         title_shape.text = slide_data['title']
#         title_shape.text_frame.paragraphs[0].font.size = Pt(36)
#         title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
#         title_shape.text_frame.paragraphs[0].font.bold = True
        
#         # Content area
#         content_left = Inches(0.5)
#         content_top = Inches(1.5)
#         content_width = Inches(9)
#         content_height = Inches(5)
        
#         # Check if slide has chart suggestion
#         if slide_data.get('chart_suggestion') and slide_data.get('chart_data'):
#             # Split slide: content on left, chart on right
#             content_width = Inches(4.5)
            
#             # Add content
#             content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#             self._add_bullet_points(content_box.text_frame, slide_data['content'])
            
#             # Add chart
#             chart_left = Inches(5)
#             self._add_chart_to_slide(slide, slide_data['chart_suggestion'], 
#                                    slide_data['chart_data'], chart_left, content_top)
#         else:
#             # Full width content
#             content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#             self._add_bullet_points(content_box.text_frame, slide_data['content'])
        
#         # Add footer
#         self._add_footer(slide)
    
#     def _add_bullet_points(self, text_frame, content_list: List[str]):
#         """Add styled bullet points to text frame"""
#         text_frame.clear()
        
#         for i, item in enumerate(content_list):
#             paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
#             paragraph.text = item
#             paragraph.font.size = Pt(18)
#             paragraph.font.color.rgb = self.text_color
#             paragraph.level = 0
#             paragraph.space_after = Pt(12)
    
#     def _add_chart_to_slide(self, slide, chart_type: str, chart_data: Dict[str, Any], 
#                            left: float, top: float):
#         """Add chart to slide"""
#         try:
#             chart_width = Inches(4)
#             chart_height = Inches(3)
            
#             if chart_type in ['bar', 'line']:
#                 # Create chart data
#                 chart_data_obj = CategoryChartData()
#                 chart_data_obj.categories = chart_data.get('labels', ['A', 'B', 'C'])
#                 chart_data_obj.add_series('Series 1', chart_data.get('values', [10, 20, 30]))
                
#                 # Add chart
#                 chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == 'bar' else XL_CHART_TYPE.LINE
#                 chart = slide.shapes.add_chart(
#                     chart_type_enum, left, top, chart_width, chart_height, chart_data_obj
#                 ).chart
                
#                 # Style chart
#                 chart.chart_title.text_frame.text = f"Data Visualization"
                
#             elif chart_type == 'pie':
#                 # For pie charts, use a text box with chart description
#                 chart_box = slide.shapes.add_textbox(left, top, chart_width, chart_height)
#                 chart_frame = chart_box.text_frame
#                 chart_frame.text = "📊 Pie Chart Data:\n"
                
#                 labels = chart_data.get('labels', ['A', 'B', 'C'])
#                 values = chart_data.get('values', [10, 20, 30])
                
#                 for label, value in zip(labels, values):
#                     chart_frame.text += f"• {label}: {value}\n"
                
#                 chart_frame.paragraphs[0].font.size = Pt(14)
#                 chart_frame.paragraphs[0].font.color.rgb = self.text_color
        
#         except Exception as e:
#             self.logger.error(f"Error adding chart: {str(e)}")
#             # Add fallback text
#             fallback_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
#             fallback_frame = fallback_box.text_frame
#             fallback_frame.text = "📊 Chart would be displayed here\n(Data visualization)"
#             fallback_frame.paragraphs[0].font.size = Pt(16)
#             fallback_frame.paragraphs[0].font.color.rgb = self.accent_color
    
#     def _add_decorative_shape(self, slide):
#         """Add decorative shape to slide"""
#         try:
#             # Add a subtle decorative rectangle
#             shape = slide.shapes.add_shape(
#                 MSO_SHAPE.RECTANGLE, 
#                 Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
#             )
#             shape.fill.solid()
#             shape.fill.fore_color.rgb = self.accent_color
#             shape.line.color.rgb = self.accent_color
#         except Exception as e:
#             self.logger.error(f"Error adding decorative shape: {str(e)}")
    
#     def _add_footer(self, slide):
#         """Add footer to slide"""
#         try:
#             footer_box = slide.shapes.add_textbox(
#                 Inches(0.5), Inches(7), Inches(9), Inches(0.5)
#             )
#             footer_frame = footer_box.text_frame
#             footer_frame.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')} | Confidential"
#             footer_frame.paragraphs[0].font.size = Pt(10)
#             footer_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
#             footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
#         except Exception as e:
#             self.logger.error(f"Error adding footer: {str(e)}")

# @st.cache_resource
# def initialize_components():
#     """Initialize all components"""
#     settings = Settings()
#     doc_processor = DocumentProcessor()
#     presentation_generator = PresentationGenerator(settings)
#     return settings, doc_processor, presentation_generator

# def main():
#     # Header
#     st.markdown("""
#     <div class="main-header">
#         <h1>🎯 Smart Presentation Generator</h1>
#         <p>Transform any content into beautiful, professional presentations with AI-powered insights and visualizations</p>
#     </div>
#     """, unsafe_allow_html=True)
    
#     # Initialize components
#     settings, doc_processor, presentation_generator = initialize_components()
    
#     # Sidebar configuration
#     with st.sidebar:
#         st.header("⚙️ Configuration")
        
#         # Input method selection
#         input_method = st.radio(
#             "Choose input method:",
#             ["📄 Upload PDF", "📝 Direct Text Input"],
#             horizontal=True
#         )
        
#         # Language selection
#         target_language = st.selectbox(
#             "Target Language",
#             options=["en", "es", "fr", "de", "it", "pt", "nl", "ru", "zh", "ja", "ko", "ar", "hi"],
#             format_func=lambda x: {
#                 "en": "🇺🇸 English", "es": "🇪🇸 Spanish", "fr": "🇫🇷 French", 
#                 "de": "🇩🇪 German", "it": "🇮🇹 Italian", "pt": "🇵🇹 Portuguese", 
#                 "nl": "🇳🇱 Dutch", "ru": "🇷🇺 Russian", "zh": "🇨🇳 Chinese", 
#                 "ja": "🇯🇵 Japanese", "ko": "🇰🇷 Korean", "ar": "🇸🇦 Arabic", 
#                 "hi": "🇮🇳 Hindi"
#             }[x]
#         )
        
#         # Company and report details
#         company_name = st.text_input("Company Name", "Your Company")
        
#         col1, col2 = st.columns(2)
#         with col1:
#             report_month = st.selectbox(
#                 "Month",
#                 options=["January", "February", "March", "April", "May", "June",
#                         "July", "August", "September", "October", "November", "December"],
#                 index=datetime.now().month - 1
#             )
#         with col2:
#             report_year = st.number_input("Year", min_value=2020, max_value=2030, value=datetime.now().year)
        
#         # Domain focus
#         domain_focus = st.multiselect(
#             "Focus Areas",
#             options=["Logistics Vehicles", "Cargo Management", "Fleet Operations", 
#                     "Supply Chain", "Transportation", "Warehouse Operations",
#                     "Business Analysis", "Performance Metrics", "Strategic Planning"],
#             default=["Business Analysis", "Performance Metrics"]
#         )
        
#         # Presentation style
#         presentation_style = st.selectbox(
#             "Presentation Style",
#             ["Professional", "Creative", "Minimalist", "Corporate"]
#         )
    
#     # Main content area
#     col1, col2 = st.columns([2, 1])
    
#     with col1:
#         if input_method == "📄 Upload PDF":
#             st.header("📄 Upload Documents")
#             uploaded_files = st.file_uploader(
#                 "Choose PDF files",
#                 type=['pdf'],
#                 accept_multiple_files=True,
#                 help="Upload one or more PDF files for analysis"
#             )
            
#             if uploaded_files:
#                 st.success(f"✅ Uploaded {len(uploaded_files)} file(s)")
                
#                 # Show file details
#                 for file in uploaded_files:
#                     st.info(f"📄 {file.name} ({file.size:,} bytes)")
        
#         else:  # Direct text input
#             st.header("📝 Text Input")
#             text_input = st.text_area(
#                 "Enter your content",
#                 placeholder="Paste your text here... (stories, reports, data, anything!)",
#                 height=200
#             )
            
#             if text_input:
#                 st.success(f"✅ Text input received ({len(text_input)} characters)")

#########code 2

# import streamlit as st
# import os
# import tempfile
# import json
# import logging
# from pathlib import Path
# from datetime import datetime
# from typing import Dict, Any, List, Optional, Union
# import zipfile
# import io
# import base64
# import re
# import boto3
# from botocore.exceptions import ClientError
# import PyPDF2
# from langdetect import detect
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.enum.shapes import MSO_SHAPE
# import pandas as pd
# import matplotlib.pyplot as plt
# import seaborn as sns
# from dotenv import load_dotenv
# import plotly.graph_objects as go
# import plotly.express as px
# from plotly.subplots import make_subplots
# import numpy as np
# from PIL import Image

# # Configure logging
# logging.basicConfig(
#     level=logging.INFO,
#     format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
#     handlers=[logging.StreamHandler(), logging.FileHandler("app.log")]
# )
# logger = logging.getLogger(__name__)

# # Custom CSS for better styling
# st.markdown("""
# <style>
#     .main-header {
#         background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
#         padding: 20px;
#         border-radius: 10px;
#         color: white;
#         margin-bottom: 30px;
#         text-align: center;
#     }
#     .feature-box {
#         background: #f8f9fa;
#         padding: 15px;
#         border-radius: 8px;
#         border-left: 4px solid #667eea;
#         margin: 10px 0;
#     }
#     .metric-container {
#         background: white;
#         padding: 20px;
#         border-radius: 10px;
#         box-shadow: 0 2px 4px rgba(0,0,0,0.1);
#         margin: 10px 0;
#     }
#     .preview-image {
#         max-width: 100%;
#         border-radius: 8px;
#         margin-top: 10px;
#     }
# </style>
# """, unsafe_allow_html=True)

# class Settings:
#     def __init__(self):
#         load_dotenv()
#         self.aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID')
#         self.aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY')
#         self.aws_session_token = os.getenv('AWS_SESSION_TOKEN')
#         self.aws_default_region = os.getenv('AWS_DEFAULT_REGION', 'us-east-1')
#         self.bedrock_model_id = os.getenv('BEDROCK_MODEL_ID', 'anthropic.claude-3-sonnet-20240229-v1:0')
#         self.bedrock_region = os.getenv('BEDROCK_REGION', 'us-east-1')
#         self.upload_folder = Path('uploads')
#         self.output_folder = Path('output')
#         self.max_file_size = 10485760  # 10MB
#         self.allowed_extensions = ['pdf']
#         self.default_source_language = 'auto'
#         self.default_target_language = 'en'
#         self.supported_languages = ['en', 'es', 'fr', 'de', 'it', 'pt', 'nl', 'ru', 'zh', 'ja', 'ko', 'ar', 'hi']
#         self.template_path = Path('templates/monthly_template.pptx')
        
#         # Create directories
#         self.upload_folder.mkdir(exist_ok=True)
#         self.output_folder.mkdir(exist_ok=True)
        
#         # Set up AWS session
#         try:
#             boto3.setup_default_session(
#                 aws_access_key_id=self.aws_access_key_id,
#                 aws_secret_access_key=self.aws_secret_access_key,
#                 aws_session_token=self.aws_session_token,
#                 region_name=self.aws_default_region
#             )
#         except Exception as e:
#             st.warning(f"AWS setup warning: {str(e)}")
#             logger.error(f"AWS setup error: {str(e)}")

# class DocumentProcessor:
#     def __init__(self):
#         self.logger = logging.getLogger(__name__)
    
#     def extract_text_from_pdf(self, pdf_path: str) -> str:
#         try:
#             text = ""
#             with open(pdf_path, 'rb') as file:
#                 pdf_reader = PyPDF2.PdfReader(file)
#                 for page in pdf_reader.pages:
#                     extracted = page.extract_text()
#                     if extracted:
#                         text += extracted + "\n"
#             return self._clean_text(text)
#         except Exception as e:
#             self.logger.error(f"Error extracting text from PDF: {str(e)}")
#             st.error(f"Error extracting PDF: {str(e)}")
#             return ""
    
#     def process_text_document(self, content: str) -> str:
#         return self._clean_text(content)
    
#     def _clean_text(self, text: str) -> str:
#         text = ' '.join(text.split())
#         text = text.replace('\x00', '').replace('\ufeff', '')
#         return text
    
#     def detect_language(self, text: str) -> Optional[str]:
#         try:
#             if len(text.strip()) < 10:
#                 return None
#             return detect(text)
#         except:
#             return None
    
#     def extract_tables(self, text: str) -> List[Dict[str, List]]:
#         """Extract table-like data from text"""
#         tables = []
#         lines = text.split('\n')
#         current_table = None
#         headers = None
        
#         for line in lines:
#             if re.search(r'[,|]\s*|\d+\s+\w+\s+\d+', line):
#                 items = [item.strip() for item in re.split(r'[,|]\s*|\s{2,}', line) if item.strip()]
#                 if len(items) >= 2:
#                     if not headers:
#                         headers = items
#                         current_table = {"headers": headers, "rows": []}
#                     else:
#                         if len(items) == len(headers):
#                             current_table["rows"].append(items)
#             else:
#                 if current_table:
#                     tables.append(current_table)
#                     headers = None
#                     current_table = None
        
#         if current_table:
#             tables.append(current_table)
        
#         return tables
    
#     def analyze_content_type(self, text: str) -> Dict[str, Any]:
#         content_analysis = {
#             'type': 'general',
#             'has_numbers': bool(re.search(r'\d+', text)),
#             'has_dates': bool(re.search(r'\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b', text, re.I)),
#             'has_tables': bool(self.extract_tables(text)),
#             'topics': [],
#             'potential_charts': []
#         }
#         text_lower = text.lower()
#         logistics_keywords = ['logistics', 'cargo', 'fleet', 'vehicle', 'transport', 'shipping', 'warehouse']
#         business_keywords = ['revenue', 'profit', 'sales', 'market', 'customer', 'performance']
#         story_keywords = ['story', 'narrative', 'character', 'plot', 'once upon']
#         sop_keywords = ['procedure', 'sop', 'standard operating', 'guideline', 'protocol', 'instruction']
        
#         if any(keyword in text_lower for keyword in logistics_keywords):
#             content_analysis['type'] = 'logistics'
#             content_analysis['topics'] = ['Fleet Operations', 'Cargo Management', 'Delivery Metrics']
#         elif any(keyword in text_lower for keyword in business_keywords):
#             content_analysis['type'] = 'business'
#             content_analysis['topics'] = ['Financials', 'Market Trends', 'Performance']
#         elif any(keyword in text_lower for keyword in story_keywords):
#             content_analysis['type'] = 'story'
#             content_analysis['topics'] = ['Narrative', 'Themes', 'Characters']
#         elif any(keyword in text_lower for keyword in sop_keywords):
#             content_analysis['type'] = 'sop'
#             content_analysis['topics'] = ['Procedures', 'Guidelines', 'Compliance']
        
#         if content_analysis['has_numbers'] or content_analysis['has_tables']:
#             content_analysis['potential_charts'] = ['bar', 'line', 'pie', 'table']
        
#         return content_analysis

# class BedrockClient:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         try:
#             self.bedrock_client = boto3.client('bedrock-runtime', region_name=self.settings.bedrock_region)
#         except Exception as e:
#             self.logger.error(f"Failed to initialize Bedrock client: {str(e)}")
#             st.error(f"Failed to initialize Bedrock client: {str(e)}")
#             self.bedrock_client = None
    
#     def invoke_model(self, prompt: str, max_tokens: int = 4000) -> Optional[str]:
#         if not self.bedrock_client:
#             st.warning("Bedrock client unavailable. Using fallback response.")
#             return self._fallback_response(prompt)
        
#         try:
#             body = {
#                 "anthropic_version": "bedrock-2023-05-31",
#                 "max_tokens": max_tokens,
#                 "messages": [{"role": "user", "content": prompt}]
#             }
#             response = self.bedrock_client.invoke_model(
#                 modelId=self.settings.bedrock_model_id,
#                 body=json.dumps(body)
#             )
#             response_body = json.loads(response['body'].read())
#             return response_body.get('content', [{}])[0].get('text', '')
#         except ClientError as e:
#             self.logger.error(f"Bedrock client error: {str(e)}")
#             st.error(f"AWS Bedrock Error: {str(e)}")
#             return self._fallback_response(prompt)
#         except Exception as e:
#             self.logger.error(f"Error invoking Bedrock model: {str(e)}")
#             st.error(f"Error invoking Bedrock: {str(e)}")
#             return self._fallback_response(prompt)
    
#     def _fallback_response(self, prompt: str) -> str:
#         if "translate" in prompt.lower():
#             return "Translation not available - using original text"
#         elif "json" in prompt.lower():
#             return json.dumps({
#                 "title": "Generated Presentation",
#                 "slides": [
#                     {
#                         "title": "Overview",
#                         "content": ["Key insights from the document", "Main topics covered", "Summary of findings"],
#                         "type": "content",
#                         "chart_suggestion": "none",
#                         "chart_data": {}
#                     }
#                 ]
#             })
#         return "Content analysis not available - using default structure"
    
#     def translate_text(self, text: str, target_language: str, source_language: str = 'auto') -> str:
#         if source_language == target_language or not text:
#             return text
#         prompt = f"""
#         Translate the following text to {target_language}, maintaining structure and formatting:
#         {text[:4000]}
#         Translation:
#         """
#         response = self.invoke_model(prompt)
#         return response or text
    
#     def generate_presentation_content(self, text: str, company_name: str, 
#                                    report_month: str, report_year: int, 
#                                    domain_focus: List[str], language: str,
#                                    content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
#         focus_areas = ", ".join(domain_focus)
#         table_summary = ""
#         if tables:
#             table_summary = "\nTables found:\n"
#             for i, table in enumerate(tables, 1):
#                 table_summary += f"Table {i}:\nHeaders: {', '.join(table['headers'])}\nRows: {len(table['rows'])}\n"
        
#         prompt = f"""
#         Analyze this {content_type} document and create a professional presentation structure for {company_name} 
#         for {report_month} {report_year} in {language}. Focus areas: {focus_areas}.
        
#         Content: {text[:4000]}
#         {table_summary}
        
#         Steps:
#         1. Identify key topics, metrics, and insights from the text and tables.
#         2. Group information into 5-7 logical slides, including a title slide.
#         3. For each slide, provide a title, 3-5 bullet points, and suggest a chart type (bar, line, pie, table, or none) with sample data.
#         4. Tailor content to the {content_type} context, using tables for data-driven slides if available.
        
#         Return JSON format:
#         {{
#             "title": "Presentation title",
#             "slides": [
#                 {{
#                     "title": "Slide title",
#                     "content": ["point 1", "point 2", "point 3"],
#                     "type": "content",
#                     "chart_suggestion": "bar|line|pie|table|none",
#                     "chart_data": {{"labels": ["A", "B"], "values": [10, 20]}} or {{"headers": [], "rows": [[]]}}
#                 }}
#             ]
#         }}
#         """
#         response = self.invoke_model(prompt, max_tokens=6000)
#         if response:
#             try:
#                 json_start = response.find('{')
#                 json_end = response.rfind('}') + 1
#                 json_str = response[json_start:json_end]
#                 return json.loads(json_str)
#             except Exception as e:
#                 self.logger.error(f"Error parsing JSON: {str(e)}")
#                 st.error(f"Error parsing Bedrock response: {str(e)}")
#         return self._create_default_structure(company_name, report_month, report_year, content_type, tables)
    
#     def _create_default_structure(self, company_name: str, report_month: str, 
#                                 report_year: int, content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
#         base_structure = {
#             "title": f"{company_name} - {report_month} {report_year} Report",
#             "slides": [
#                 {
#                     "title": "Executive Summary",
#                     "content": ["Overview of key findings", "Main topics identified", "Summary of performance"],
#                     "type": "content",
#                     "chart_suggestion": "none",
#                     "chart_data": {}
#                 },
#                 {
#                     "title": "Key Insights",
#                     "content": ["Primary observations", "Notable trends", "Actionable takeaways"],
#                     "type": "content",
#                     "chart_suggestion": "none",
#                     "chart_data": {}
#                 }
#             ]
#         }
#         if content_type == 'logistics':
#             base_structure["slides"].append({
#                 "title": "Logistics Metrics",
#                 "content": ["Operational efficiency", "Fleet performance", "Delivery metrics"],
#                 "type": "content",
#                 "chart_suggestion": "bar",
#                 "chart_data": {"labels": ["Efficiency", "On-Time", "Cost"], "values": [85, 92, 75]}
#             })
#         elif content_type == 'business':
#             base_structure["slides"].append({
#                 "title": "Financial Overview",
#                 "content": ["Revenue trends", "Profit margins", "Cost analysis"],
#                 "type": "content",
#                 "chart_suggestion": "line",
#                 "chart_data": {"labels": ["Q1", "Q2", "Q3"], "values": [100, 120, 130]}
#             })
#         elif content_type == 'story':
#             base_structure["slides"].append({
#                 "title": "Key Themes",
#                 "content": ["Main narrative themes", "Character development", "Plot highlights"],
#                 "type": "content",
#                 "chart_suggestion": "none",
#                 "chart_data": {}
#             })
#         elif content_type == 'sop':
#             base_structure["slides"].append({
#                 "title": "Procedure Overview",
#                 "content": ["Key procedures", "Compliance requirements", "Implementation steps"],
#                 "type": "content",
#                 "chart_suggestion": "table",
#                 "chart_data": {"headers": ["Step", "Description"], "rows": [["1", "Initiate process"], ["2", "Verify compliance"]]}
#             })
#         if tables:
#             base_structure["slides"].append({
#                 "title": "Data Summary",
#                 "content": [f"Table with {len(tables[0]['headers'])} columns and {len(tables[0]['rows'])} rows"],
#                 "type": "content",
#                 "chart_suggestion": "table",
#                 "chart_data": {"headers": tables[0]["headers"], "rows": tables[0]["rows"][:3]}
#             })
#         return base_structure

# class ChartGenerator:
#     def __init__(self):
#         self.colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
    
#     def create_chart(self, chart_type: str, data: Dict[str, Any], style: str) -> str:
#         """Create chart using Plotly and return base64 encoded image"""
#         try:
#             fig = None
#             labels = data.get('labels', ['A', 'B', 'C'])
#             values = data.get('values', [10, 20, 30])
#             headers = data.get('headers', [])
#             rows = data.get('rows', [])
            
#             plt.style.use('seaborn-v0_8')
#             if chart_type == 'bar':
#                 fig = px.bar(x=labels, y=values, color=labels, color_discrete_sequence=self.colors)
#                 fig.update_layout(title="Data Visualization", showlegend=False)
#             elif chart_type == 'line':
#                 fig = px.line(x=labels, y=values, markers=True, color_discrete_sequence=[self.colors[0]])
#                 fig.update_layout(title="Trend Analysis")
#             elif chart_type == 'pie':
#                 fig = px.pie(values=values, names=labels, color_discrete_sequence=self.colors)
#                 fig.update_layout(title="Distribution")
#             elif chart_type == 'table':
#                 if headers and rows:
#                     fig = go.Figure(data=[go.Table(
#                         header=dict(values=headers, fill_color=self.colors[0], align='center'),
#                         cells=dict(values=[rows[i] for i in range(len(rows))], fill_color='lavender', align='left')
#                     )])
#                     fig.update_layout(title="Data Table")
            
#             if fig:
#                 if style == "Minimalist":
#                     fig.update_layout(showlegend=False, plot_bgcolor='rgba(0,0,0,0)', paper_bgcolor='rgba(0,0,0,0)')
#                 elif style == "Creative":
#                     fig.update_layout(font=dict(size=14, color='#333333'), title_font_size=20)
                
#                 buffer = io.BytesIO()
#                 fig.write_image(buffer, format='png', scale=2)
#                 buffer.seek(0)
#                 chart_base64 = base64.b64encode(buffer.read()).decode()
#                 return chart_base64
#             return ""
#         except Exception as e:
#             logger.error(f"Error creating chart: {str(e)}")
#             return ""

# class PresentationGenerator:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         self.chart_generator = ChartGenerator()
#         self.theme = {
#             "colors": {
#                 "primary": "#003087",
#                 "secondary": "#FFFFFF",
#                 "accent": "#FFA500",
#                 "text": "#404040",
#                 "footer": "#808080"
#             },
#             "fonts": {
#                 "title": {"name": "Calibri", "size": 44, "bold": True},
#                 "subtitle": {"name": "Calibri", "size": 24, "bold": False},
#                 "content": {"name": "Calibri", "size": 18, "bold": False},
#                 "footer": {"name": "Calibri", "size": 10, "bold": False}
#             },
#             "footer": {
#                 "text": "{company_name} | {date} | Confidential",
#                 "position": {"left": 0.5, "top": 7.0, "width": 9.0, "height": 0.5}
#             }
#         }
#         try:
#             theme_path = Path("templates/styles/corporate_theme.json")
#             if theme_path.exists():
#                 with open(theme_path, "r") as f:
#                     self.theme = json.load(f)["theme"]
#         except Exception as e:
#             self.logger.warning(f"Failed to load theme: {str(e)}. Using default theme.")
#         self.primary_color = RGBColor.from_string(self.theme["colors"]["primary"].lstrip("#"))
#         self.secondary_color = RGBColor.from_string(self.theme["colors"]["secondary"].lstrip("#"))
#         self.accent_color = RGBColor.from_string(self.theme["colors"]["accent"].lstrip("#"))
    
#     def generate_presentation(self, content: str, company_name: str, 
#                            report_month: str, report_year: int, 
#                            domain_focus: List[str], language: str,
#                            content_type: str, tables: List[Dict[str, List]], style: str) -> str:
#         try:
#             # Try to load template, fall back to default if not found
#             try:
#                 prs = Presentation(str(self.settings.template_path))
#             except Exception as e:
#                 self.logger.warning(f"Template not found at {self.settings.template_path}: {str(e)}. Using default presentation.")
#                 st.warning(f"Using default presentation template as {self.settings.template_path} is missing.")
#                 prs = Presentation()
            
#             bedrock_client = BedrockClient(self.settings)
#             content_structure = bedrock_client.generate_presentation_content(
#                 content, company_name, report_month, report_year, domain_focus, language, content_type, tables
#             )
#             if not content_structure:
#                 content_structure = bedrock_client._create_default_structure(company_name, report_month, report_year, content_type, tables)
            
#             self._create_enhanced_title_slide(prs, content_structure['title'], company_name, report_month, report_year, domain_focus, style)
#             for slide_data in content_structure['slides']:
#                 self._create_enhanced_content_slide(prs, slide_data, style)
            
#             filename = f"Presentation_{company_name.replace(' ', '_')}_{report_month}_{report_year}.pptx"
#             output_path = self.settings.output_folder / filename
#             prs.save(str(output_path))
#             return str(output_path)
#         except Exception as e:
#             self.logger.error(f"Error generating presentation: {str(e)}")
#             st.error(f"Error generating presentation: {str(e)}")
#             raise
    
#     def _create_enhanced_title_slide(self, prs: Presentation, title: str, company_name: str, 
#                                   report_month: str, report_year: int, domain_focus: List[str], style: str):
#         slide_layout = prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
        
#         for shape in slide.shapes:
#             if hasattr(shape, 'text'):
#                 shape.text = ""
        
#         left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(2)
#         title_box = slide.shapes.add_textbox(left, top, width, height)
#         title_frame = title_box.text_frame
#         title_frame.text = title
#         title_paragraph = title_frame.paragraphs[0]
#         title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
#         title_paragraph.font.color.rgb = self.primary_color
#         title_paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
#         title_paragraph.alignment = PP_ALIGN.CENTER
        
#         subtitle_top = Inches(3.5)
#         subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1.5))
#         subtitle_frame = subtitle_box.text_frame
#         subtitle_frame.text = f"{company_name} | {report_month} {report_year} | {', '.join(domain_focus)}"
#         subtitle_paragraph = subtitle_frame.paragraphs[0]
#         subtitle_paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
#         subtitle_paragraph.font.color.rgb = self.accent_color
#         subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
#         if style in ["Creative", "Minimalist"]:
#             self._add_decorative_image(slide, style)
#         self._add_footer(slide, company_name, f"{report_month} {report_year}")
    
#     def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], style: str):
#         slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
        
#         title_shape = slide.shapes.title
#         title_shape.text = slide_data['title']
#         title_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["title"]["size"])
#         title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
#         title_shape.text_frame.paragraphs[0].font.bold = True
        
#         content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
#         if slide_data.get('chart_suggestion') != "none" and slide_data.get('chart_data'):
#             content_width = Inches(4.5)
#             content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#             self._add_bullet_points(content_box.text_frame, slide_data['content'])
#             chart_left = Inches(5.0)
#             self._add_chart_to_slide(slide, slide_data['chart_suggestion'], slide_data['chart_data'], chart_left, content_top, style)
#         else:
#             content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#             self._add_bullet_points(content_box.text_frame, slide_data['content'])
        
#         self._add_footer(slide, "Your Company", f"{datetime.now().strftime('%B %Y')}")
    
#     def _add_bullet_points(self, text_frame, content_list: List[str]):
#         text_frame.clear()
#         for i, item in enumerate(content_list):
#             paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
#             paragraph.text = item
#             paragraph.font.size = Pt(self.theme["fonts"]["content"]["size"])
#             paragraph.font.color.rgb = RGBColor.from_string(self.theme["colors"]["text"].lstrip("#"))
#             paragraph.level = 0
#             paragraph.space_after = Pt(12)
    
#     def _add_chart_to_slide(self, slide, chart_type: str, chart_data: Dict[str, Any], left: float, top: float, style: str):
#         try:
#             chart_width, chart_height = Inches(4.5), Inches(3.5)
#             if chart_type == 'table':
#                 table = slide.shapes.add_table(
#                     len(chart_data.get('rows', [])[:5]) + 1,
#                     len(chart_data.get('headers', [])),
#                     left, top, chart_width, chart_height
#                 ).table
#                 for i, header in enumerate(chart_data.get('headers', [])):
#                     table.cell(0, i).text = header
#                     table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)
#                 for i, row in enumerate(chart_data.get('rows', [])[:5], 1):
#                     for j, value in enumerate(row):
#                         table.cell(i, j).text = str(value)
#                         table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
#             else:
#                 chart_base64 = self.chart_generator.create_chart(chart_type, chart_data, style)
#                 if chart_base64:
#                     buffer = io.BytesIO(base64.b64decode(chart_base64))
#                     slide.shapes.add_picture(buffer, left, top, chart_width, chart_height)
#         except Exception as e:
#             self.logger.error(f"Error adding chart: {str(e)}")
#             st.error(f"Error adding chart: {str(e)}")
#             fallback_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
#             fallback_box.text_frame.text = "📊 Chart unavailable\n(Data visualization)"
#             fallback_box.text_frame.paragraphs[0].font.size = Pt(16)
#             fallback_box.text_frame.paragraphs[0].font.color.rgb = self.accent_color
    
#     def _add_decorative_image(self, slide, style: str):
#         """Add a decorative image or icon based on style"""
#         try:
#             img_path = Path("templates/icons/placeholder.png")
#             if not img_path.exists():
#                 img = Image.new('RGB', (200, 200), color=(255, 255, 255))
#                 img.save(img_path)
#             left, top = Inches(8.5), Inches(0.5)
#             if style == "Creative":
#                 slide.shapes.add_picture(str(img_path), left, top, width=Inches(1.5))
#             elif style == "Minimalist":
#                 slide.shapes.add_picture(str(img_path), left, top, width=Inches(1))
#         except Exception as e:
#             self.logger.error(f"Error adding decorative image: {str(e)}")
    
#     def _add_footer(self, slide, company_name: str, date: str):
#         footer_box = slide.shapes.add_textbox(
#             Inches(self.theme["footer"]["position"]["left"]),
#             Inches(self.theme["footer"]["position"]["top"]),
#             Inches(self.theme["footer"]["position"]["width"]),
#             Inches(self.theme["footer"]["position"]["height"])
#         )
#         footer_frame = footer_box.text_frame
#         footer_frame.text = self.theme["footer"]["text"].format(company_name=company_name, date=date)
#         footer_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
#         footer_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(self.theme["colors"]["footer"].lstrip("#"))
#         footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

# @st.cache_resource
# def initialize_components():
#     settings = Settings()
#     doc_processor = DocumentProcessor()
#     presentation_generator = PresentationGenerator(settings)
#     return settings, doc_processor, presentation_generator

# def main():
#     st.markdown("""
#     <div class="main-header">
#         <h1>🎯 Smart Presentation Generator</h1>
#         <p>Transform any content into beautiful, professional presentations with AI-powered insights and visualizations</p>
#     </div>
#     """, unsafe_allow_html=True)
    
#     settings, doc_processor, presentation_generator = initialize_components()
#     content_type = "Not yet processed"
    
#     with st.sidebar:
#         st.header("⚙️ Configuration")
#         input_method = st.radio(
#             "Choose input method:",
#             ["📄 Upload PDF", "📝 Direct Text Input"],
#             horizontal=True
#         )
#         target_language = st.selectbox(
#             "Target Language",
#             options=["en", "es", "fr", "de", "it", "pt", "nl", "ru", "zh", "ja", "ko", "ar", "hi"],
#             format_func=lambda x: {
#                 "en": "🇺🇸 English", "es": "🇪🇸 Spanish", "fr": "🇫🇷 French", 
#                 "de": "🇩🇪 German", "it": "🇮🇹 Italian", "pt": "🇵🇹 Portuguese", 
#                 "nl": "🇳🇱 Dutch", "ru": "🇷🇺 Russian", "zh": "🇨🇳 Chinese", 
#                 "ja": "🇯🇵 Japanese", "ko": "🇰🇷 Korean", "ar": "🇸🇦 Arabic", 
#                 "hi": "🇮🇳 Hindi"
#             }[x]
#         )
#         company_name = st.text_input("Company Name", "Your Company")
#         col1, col2 = st.columns(2)
#         with col1:
#             report_month = st.selectbox(
#                 "Month",
#                 options=["January", "February", "March", "April", "May", "June",
#                         "July", "August", "September", "October", "November", "December"],
#                 index=datetime.now().month - 1
#             )
#         with col2:
#             report_year = st.number_input("Year", min_value=2020, max_value=2030, value=datetime.now().year)
#         domain_focus = st.multiselect(
#             "Focus Areas",
#             options=["Logistics Vehicles", "Cargo Management", "Fleet Operations", 
#                     "Supply Chain", "Transportation", "Warehouse Operations",
#                     "Business Analysis", "Performance Metrics", "Strategic Planning"],
#             default=["Business Analysis", "Performance Metrics"]
#         )
#         presentation_style = st.selectbox(
#             "Presentation Style",
#             ["Professional", "Creative", "Minimalist", "Corporate"]
#         )
    
#     col1, col2 = st.columns([2, 1])
    
#     with col1:
#         if input_method == "📄 Upload PDF":
#             st.header("📄 Upload Documents")
#             uploaded_files = st.file_uploader(
#                 "Choose PDF files",
#                 type=['pdf'],
#                 accept_multiple_files=True,
#                 help="Upload one or more PDF files for analysis"
#             )
#             if uploaded_files:
#                 st.success(f"✅ Uploaded {len(uploaded_files)} file(s)")
#                 for file in uploaded_files:
#                     st.info(f"📄 {file.name} ({file.size:,} bytes)")
#         else:
#             st.header("📝 Text Input")
#             text_input = st.text_area(
#                 "Enter your content",
#                 placeholder="Paste your text here... (stories, reports, data, anything!)",
#                 height=200
#             )
#             if text_input:
#                 st.success(f"✅ Text input received ({len(text_input)} characters)")
        
#         if (input_method == "📄 Upload PDF" and uploaded_files) or (input_method == "📝 Direct Text Input" and text_input):
#             if st.button("🚀 Generate Presentation", type="primary"):
#                 with st.spinner("Processing content and generating presentation..."):
#                     try:
#                         progress_bar = st.progress(0)
#                         status_text = st.empty()
                        
#                         status_text.text("Processing input...")
#                         all_text = ""
#                         tables = []
#                         if input_method == "📄 Upload PDF":
#                             for i, file in enumerate(uploaded_files):
#                                 with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
#                                     tmp_file.write(file.read())
#                                     tmp_file_path = tmp_file.name
#                                 text = doc_processor.extract_text_from_pdf(tmp_file_path)
#                                 all_text += f"\n\n--- {file.name} ---\n{text}"
#                                 tables.extend(doc_processor.extract_tables(text))
#                                 os.unlink(tmp_file_path)
#                                 progress_bar.progress((i + 1) / len(uploaded_files) * 0.3)
#                         else:
#                             all_text = doc_processor.process_text_document(text_input)
#                             tables = doc_processor.extract_tables(all_text)
#                             progress_bar.progress(0.3)
                        
#                         status_text.text("Analyzing content...")
#                         content_analysis = doc_processor.analyze_content_type(all_text)
#                         content_type = content_analysis['type']
#                         progress_bar.progress(0.5)
                        
#                         status_text.text("Translating content...")
#                         bedrock_client = BedrockClient(settings)
#                         source_lang = doc_processor.detect_language(all_text)
#                         translated_text = bedrock_client.translate_text(all_text, target_language, source_lang or 'auto')
#                         progress_bar.progress(0.7)
                        
#                         status_text.text("Generating presentation...")
#                         presentation_path = presentation_generator.generate_presentation(
#                             translated_text,
#                             company_name,
#                             report_month,
#                             report_year,
#                             domain_focus,
#                             target_language,
#                             content_type,
#                             tables,
#                             presentation_style
#                         )
#                         progress_bar.progress(1.0)
                        
#                         status_text.text("Presentation generated successfully!")
#                         st.success("✅ Presentation generated successfully!")
#                         with open(presentation_path, 'rb') as file:
#                             st.download_button(
#                                 label="📥 Download Presentation",
#                                 data=file.read(),
#                                 file_name=f"Presentation_{report_month}_{report_year}.pptx",
#                                 mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#                             )
                        
#                         content_structure = bedrock_client.generate_presentation_content(
#                             translated_text, company_name, report_month, report_year, 
#                             domain_focus, target_language, content_type, tables
#                         )
#                         for slide_data in content_structure['slides']:
#                             if slide_data.get('chart_suggestion') != "none" and slide_data.get('chart_data'):
#                                 chart_base64 = presentation_generator.chart_generator.create_chart(
#                                     slide_data['chart_suggestion'], slide_data['chart_data'], presentation_style
#                                 )
#                                 if chart_base64:
#                                     st.image(f"data:image/png;base64,{chart_base64}", caption="Sample Chart Preview", use_column_width=True)
#                                     break
                        
#                     except ClientError as e:
#                         st.error(f"AWS Bedrock Error: {str(e)}")
#                         st.warning("Generated default presentation due to AWS error. Please check your AWS credentials.")
#                     except Exception as e:
#                         st.error(f"Error: {str(e)}")
#                         st.warning("Generated default presentation due to processing error.")
    
#     with col2:
#         st.header("📋 Preview")
#         if input_method == "📄 Upload PDF" and uploaded_files:
#             st.subheader("Uploaded Files")
#             for file in uploaded_files:
#                 st.write(f"• {file.name}")
#         elif input_method == "📝 Direct Text Input" and text_input:
#             st.subheader("Text Input Preview")
#             st.write(f"First {min(100, len(text_input))} characters: {text_input[:100]}...")
        
#         st.subheader("Output Configuration")
#         st.write(f"**Input Method:** {input_method}")
#         st.write(f"**Target Language:** {target_language}")
#         st.write(f"**Company:** {company_name}")
#         st.write(f"**Period:** {report_month} {report_year}")
#         st.write(f"**Focus Areas:** {', '.join(domain_focus)}")
#         st.write(f"**Content Type:** {content_type}")
#         st.write(f"**Presentation Style:** {presentation_style}")
        
#         st.subheader("📊 Benefits")
#         with st.container():
#             st.markdown('<div class="metric-container">', unsafe_allow_html=True)
#             st.metric("Time Saved", "50%", "vs manual creation")
#             st.metric("Cost Reduction", "70%", "vs professional services")
#             st.metric("GDPR Compliance", "100%", "Data stays in AWS")
#             st.markdown('</div>', unsafe_allow_html=True)

# if __name__ == "__main__":
#     main()


########## code 3
# import streamlit as st
# import os
# import tempfile
# import json
# import logging
# from pathlib import Path
# from datetime import datetime
# from typing import Dict, Any, List, Optional, Union
# import zipfile
# import io
# import base64
# import re
# import boto3
# from botocore.exceptions import ClientError
# import PyPDF2
# from langdetect import detect
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
# from pptx.dml.color import RGBColor
# from pptx.chart.data import CategoryChartData
# from pptx.enum.chart import XL_CHART_TYPE
# from pptx.enum.shapes import MSO_SHAPE
# import pandas as pd
# import matplotlib.pyplot as plt
# import seaborn as sns
# from dotenv import load_dotenv
# import plotly.graph_objects as go
# import plotly.express as px
# from plotly.subplots import make_subplots
# import numpy as np
# from PIL import Image

# # Configure logging
# logging.basicConfig(
#     level=logging.INFO,
#     format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
#     handlers=[logging.StreamHandler(), logging.FileHandler("app.log")]
# )
# logger = logging.getLogger(__name__)

# # Custom CSS for better styling
# st.markdown("""
# <style>
#     .main-header {
#         background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
#         padding: 20px;
#         border-radius: 10px;
#         color: white;
#         margin-bottom: 30px;
#         text-align: center;
#     }
#     .feature-box {
#         background: #f8f9fa;
#         padding: 15px;
#         border-radius: 8px;
#         border-left: 4px solid #667eea;
#         margin: 10px 0;
#     }
#     .metric-container {
#         background: white;
#         padding: 20px;
#         border-radius: 10px;
#         box-shadow: 0 2px 4px rgba(0,0,0,0.1);
#         margin: 10px 0;
#     }
#     .preview-image {
#         max-width: 100%;
#         border-radius: 8px;
#         margin-top: 10px;
#     }
# </style>
# """, unsafe_allow_html=True)

# class Settings:
#     def __init__(self):
#         load_dotenv()
#         self.aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID')
#         self.aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY')
#         self.aws_session_token = os.getenv('AWS_SESSION_TOKEN')
#         self.aws_default_region = os.getenv('AWS_DEFAULT_REGION', 'us-east-1')
#         self.bedrock_model_id = os.getenv('BEDROCK_MODEL_ID', 'anthropic.claude-3-sonnet-20240229-v1:0')
#         self.bedrock_region = os.getenv('BEDROCK_REGION', 'us-east-1')
        
#         # Robust template path resolution
#         self.template_path = Path(__file__).parent / 'templates' / 'monthly_template.pptx'
        
#         self.upload_folder = Path('uploads')
#         self.output_folder = Path('output')
#         self.max_file_size = 10485760  # 10MB
#         self.allowed_extensions = ['pdf']
#         self.default_source_language = 'auto'
#         self.default_target_language = 'en'
#         self.supported_languages = ['en', 'es', 'fr', 'de', 'it', 'pt', 'nl', 'ru', 'zh', 'ja', 'ko', 'ar', 'hi']
        
#         # Create directories, including templates/styles if they don't exist
#         self.upload_folder.mkdir(exist_ok=True)
#         self.output_folder.mkdir(exist_ok=True)
#         (Path(__file__).parent / 'templates').mkdir(exist_ok=True)
#         (Path(__file__).parent / 'templates' / 'styles').mkdir(exist_ok=True)
#         (Path(__file__).parent / 'templates' / 'icons').mkdir(exist_ok=True) # Ensure icons folder exists

#         # Set up AWS session
#         try:
#             boto3.setup_default_session(
#                 aws_access_key_id=self.aws_access_key_id,
#                 aws_secret_access_key=self.aws_secret_access_key,
#                 aws_session_token=self.aws_session_token,
#                 region_name=self.aws_default_region
#             )
#         except Exception as e:
#             st.warning(f"AWS setup warning: {str(e)}")
#             logger.error(f"AWS setup error: {str(e)}")

# class DocumentProcessor:
#     def __init__(self):
#         self.logger = logging.getLogger(__name__)
    
#     def extract_text_from_pdf(self, pdf_path: str) -> str:
#         try:
#             text = ""
#             with open(pdf_path, 'rb') as file:
#                 pdf_reader = PyPDF2.PdfReader(file)
#                 for page in pdf_reader.pages:
#                     extracted = page.extract_text()
#                     if extracted:
#                         text += extracted + "\n"
#             return self._clean_text(text)
#         except Exception as e:
#             self.logger.error(f"Error extracting text from PDF: {str(e)}")
#             st.error(f"Error extracting PDF: {str(e)}")
#             return ""
    
#     def process_text_document(self, content: str) -> str:
#         return self._clean_text(content)
    
#     def _clean_text(self, text: str) -> str:
#         text = ' '.join(text.split())
#         text = text.replace('\x00', '').replace('\ufeff', '')
#         return text
    
#     def detect_language(self, text: str) -> Optional[str]:
#         try:
#             if len(text.strip()) < 10:
#                 return None
#             return detect(text)
#         except:
#             return None
    
#     def extract_tables(self, text: str) -> List[Dict[str, List]]:
#         """Extract table-like data from text"""
#         tables = []
#         lines = text.split('\n')
#         current_table = None
#         headers = None
        
#         for line in lines:
#             if re.search(r'[,|]\s*|\d+\s+\w+\s+\d+', line):
#                 items = [item.strip() for item in re.split(r'[,|]\s*|\s{2,}', line) if item.strip()]
#                 if len(items) >= 2:
#                     if not headers:
#                         headers = items
#                         current_table = {"headers": headers, "rows": []}
#                     else:
#                         if len(items) == len(headers):
#                             current_table["rows"].append(items)
#             else:
#                 if current_table:
#                     tables.append(current_table)
#                     headers = None
#                     current_table = None
        
#         if current_table:
#             tables.append(current_table)
        
#         return tables
    
#     def analyze_content_type(self, text: str) -> Dict[str, Any]:
#         content_analysis = {
#             'type': 'general',
#             'has_numbers': bool(re.search(r'\d+', text)),
#             'has_dates': bool(re.search(r'\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b', text, re.I)),
#             'has_tables': bool(self.extract_tables(text)),
#             'topics': [],
#             'potential_charts': []
#         }
#         text_lower = text.lower()
#         logistics_keywords = ['logistics', 'cargo', 'fleet', 'vehicle', 'transport', 'shipping', 'warehouse']
#         business_keywords = ['revenue', 'profit', 'sales', 'market', 'customer', 'performance']
#         story_keywords = ['story', 'narrative', 'character', 'plot', 'once upon']
#         sop_keywords = ['procedure', 'sop', 'standard operating', 'guideline', 'protocol', 'instruction']
        
#         if any(keyword in text_lower for keyword in logistics_keywords):
#             content_analysis['type'] = 'logistics'
#             content_analysis['topics'] = ['Fleet Operations', 'Cargo Management', 'Delivery Metrics']
#         elif any(keyword in text_lower for keyword in business_keywords):
#             content_analysis['type'] = 'business'
#             content_analysis['topics'] = ['Financials', 'Market Trends', 'Performance']
#         elif any(keyword in text_lower for keyword in story_keywords):
#             content_analysis['type'] = 'story'
#             content_analysis['topics'] = ['Narrative', 'Themes', 'Characters']
#         elif any(keyword in text_lower for keyword in sop_keywords):
#             content_analysis['type'] = 'sop'
#             content_analysis['topics'] = ['Procedures', 'Guidelines', 'Compliance']
        
#         if content_analysis['has_numbers'] or content_analysis['has_tables']:
#             content_analysis['potential_charts'] = ['bar', 'line', 'pie', 'table']
        
#         # Add statistical analysis
#         statistical_data = self.extract_statistical_data(text)
#         content_analysis['statistical_data'] = statistical_data
#         content_analysis['has_statistics'] = len(statistical_data['numbers']) > 0 or len(statistical_data['percentages']) > 0

#         # Suggest chart types based on data
#         if statistical_data['percentages']:
#             content_analysis['potential_charts'].extend(['pie', 'donut'])
#         if statistical_data['comparisons']:
#             content_analysis['potential_charts'].extend(['bar', 'column'])
#         if statistical_data['monetary_values']:
#             content_analysis['potential_charts'].extend(['line', 'area'])
#         return content_analysis

#     def extract_statistical_data(self, text: str) -> Dict[str, Any]:
#         """Extract numerical data, percentages, and statistical information from text"""
#         statistical_data = {
#             'numbers': [],
#             'percentages': [],
#             'dates': [],
#             'monetary_values': [],
#             'metrics': [],
#             'comparisons': []
#         }
        
#         # Extract numbers with context
#         number_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(%|percent|million|billion|thousand|units|kg|tons|miles|km)?'
#         numbers = re.findall(number_pattern, text, re.IGNORECASE)
        
#         for context, value, unit in numbers:
#             try:
#                 numeric_value = float(value.replace(',', ''))
#                 statistical_data['numbers'].append({
#                     'context': context.strip(),
#                     'value': numeric_value,
#                     'unit': unit,
#                     'original': f"{value} {unit}".strip()
#                 })
#             except ValueError:
#                 continue
        
#         # Extract percentages
#         percentage_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:\.\d+)?)\s*(%|percent)'
#         percentages = re.findall(percentage_pattern, text, re.IGNORECASE)
        
#         for context, value, unit in percentages:
#             try:
#                 statistical_data['percentages'].append({
#                     'context': context.strip(),
#                     'value': float(value),
#                     'original': f"{value}%"
#                 })
#             except ValueError:
#                 continue
        
#         # Extract monetary values
#         money_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*\$(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|thousand|M|B|K)?'
#         monetary = re.findall(money_pattern, text, re.IGNORECASE)
        
#         for context, value, multiplier in monetary:
#             try:
#                 numeric_value = float(value.replace(',', ''))
#                 multiplier_map = {'million': 1000000, 'M': 1000000, 'billion': 1000000000, 'B': 1000000000, 'thousand': 1000, 'K': 1000}
#                 if multiplier and multiplier in multiplier_map:
#                     numeric_value *= multiplier_map[multiplier]
                
#                 statistical_data['monetary_values'].append({
#                     'context': context.strip(),
#                     'value': numeric_value,
#                     'original': f"${value} {multiplier}".strip()
#                 })
#             except ValueError:
#                 continue
        
#         # Extract comparison data (increased by X%, decreased by Y%)
#         comparison_pattern = r'(\w+(?:\s+\w+){0,2})\s*(?:increased|decreased|rose|fell|grew|declined)\s*(?:by\s*)?(\d+(?:\.\d+)?)\s*(%|percent)'
#         comparisons = re.findall(comparison_pattern, text, re.IGNORECASE)
        
#         for context, value, unit in comparisons:
#             try:
#                 statistical_data['comparisons'].append({
#                     'context': context.strip(),
#                     'value': float(value),
#                     'type': 'increase' if any(word in text.lower() for word in ['increased', 'rose', 'grew']) else 'decrease'
#                 })
#             except ValueError:
#                 continue
        
#         return statistical_data

# class BedrockClient:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         try:
#             self.bedrock_client = boto3.client('bedrock-runtime', region_name=self.settings.bedrock_region)
#         except Exception as e:
#             self.logger.error(f"Failed to initialize Bedrock client: {str(e)}")
#             st.error(f"Failed to initialize Bedrock client: {str(e)}")
#             self.bedrock_client = None
    
#     def invoke_model(self, prompt: str, max_tokens: int = 4000) -> Optional[str]:
#         if not self.bedrock_client:
#             st.warning("Bedrock client unavailable. Using fallback response.")
#             return self._fallback_response(prompt)
        
#         try:
#             body = {
#                 "anthropic_version": "bedrock-2023-05-31",
#                 "max_tokens": max_tokens,
#                 "messages": [{"role": "user", "content": prompt}]
#             }
#             response = self.bedrock_client.invoke_model(
#                 modelId=self.settings.bedrock_model_id,
#                 body=json.dumps(body)
#             )
#             response_body = json.loads(response['body'].read())
#             return response_body.get('content', [{}])[0].get('text', '')
#         except ClientError as e:
#             self.logger.error(f"Bedrock client error: {str(e)}")
#             st.error(f"AWS Bedrock Error: {str(e)}")
#             return self._fallback_response(prompt)
#         except Exception as e:
#             self.logger.error(f"Error invoking Bedrock model: {str(e)}")
#             st.error(f"Error invoking Bedrock: {str(e)}")
#             return self._fallback_response(prompt)
    
#     def _fallback_response(self, prompt: str) -> str:
#         if "translate" in prompt.lower():
#             return "Translation not available - using original text"
#         elif "json" in prompt.lower():
#             return json.dumps({
#                 "title": "Generated Presentation",
#                 "slides": [
#                     {
#                         "title": "Overview",
#                         "content": ["Key insights from the document", "Main topics covered", "Summary of findings"],
#                         "type": "content",
#                         "chart_suggestion": "none",
#                         "chart_data": {}
#                     }
#                 ]
#             })
#         return "Content analysis not available - using default structure"
    
#     def translate_text(self, text: str, target_language: str, source_language: str = 'auto') -> str:
#         if source_language == target_language or not text:
#             return text
#         prompt = f"""
#         Translate the following text to {target_language}, maintaining structure and formatting:
#         {text[:4000]}
#         Translation:
#         """
#         response = self.invoke_model(prompt)
#         return response or text
    
#     def generate_presentation_content(self, text: str, company_name: str, 
#                                       report_month: str, report_year: int, 
#                                       domain_focus: List[str], language: str,
#                                       content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
#         focus_areas = ", ".join(domain_focus)
#         table_summary = ""
#         if tables:
#             table_summary = "\nTables found:\n"
#             for i, table in enumerate(tables, 1):
#                 table_summary += f"Table {i}:\nHeaders: {', '.join(table['headers'])}\nRows: {len(table['rows'])}\n"
        
#         prompt = f"""
#         Analyze this {content_type} document and create a professional presentation structure for {company_name} 
#         for {report_month} {report_year} in {language}. Focus areas: {focus_areas}.
        
#         Content: {text[:4000]}
#         {table_summary}
        
#         Steps:
#         1. Identify key topics, metrics, and insights from the text and tables.
#         2. Group information into 5-7 logical slides, including a title slide.
#         3. For each slide, provide a title, 3-5 bullet points, and suggest a chart type (bar, line, pie, table, or none) with sample data.
#         4. Tailor content to the {content_type} context, using tables for data-driven slides if available.
        
#         Return JSON format:
#         {{
#             "title": "Presentation title",
#             "slides": [
#                 {{
#                     "title": "Slide title",
#                     "content": ["point 1", "point 2", "point 3"],
#                     "type": "content",
#                     "chart_suggestion": "bar|line|pie|table|none",
#                     "chart_data": {{"labels": ["A", "B"], "values": [10, 20]}} or {{"headers": [], "rows": [[]]}}
#                 }}
#             ]
#         }}
#         """
#         response = self.invoke_model(prompt, max_tokens=6000)
#         if response:
#             try:
#                 json_start = response.find('{')
#                 json_end = response.rfind('}') + 1
#                 json_str = response[json_start:json_end]
#                 return json.loads(json_str)
#             except Exception as e:
#                 self.logger.error(f"Error parsing JSON: {str(e)}")
#                 st.error(f"Error parsing Bedrock response: {str(e)}")
#         return self._create_default_structure(company_name, report_month, report_year, content_type, tables)
    
#     def _create_default_structure(self, company_name: str, report_month: str, 
#                                  report_year: int, content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
#         base_structure = {
#             "title": f"{company_name} - {report_month} {report_year} Report",
#             "slides": [
#                 {
#                     "title": "Executive Summary",
#                     "content": ["Overview of key findings", "Main topics identified", "Summary of performance"],
#                     "type": "content",
#                     "chart_suggestion": "none",
#                     "chart_data": {}
#                 },
#                 {
#                     "title": "Key Insights",
#                     "content": ["Primary observations", "Notable trends", "Actionable takeaways"],
#                     "type": "content",
#                     "chart_suggestion": "none",
#                     "chart_data": {}
#                 }
#             ]
#         }
#         if content_type == 'logistics':
#             base_structure["slides"].append({
#                 "title": "Logistics Metrics",
#                 "content": ["Operational efficiency", "Fleet performance", "Delivery metrics"],
#                 "type": "content",
#                 "chart_suggestion": "bar",
#                 "chart_data": {"labels": ["Efficiency", "On-Time", "Cost"], "values": [85, 92, 75]}
#             })
#         elif content_type == 'business':
#             base_structure["slides"].append({
#                 "title": "Financial Overview",
#                 "content": ["Revenue trends", "Profit margins", "Cost analysis"],
#                 "type": "content",
#                 "chart_suggestion": "line",
#                 "chart_data": {"labels": ["Q1", "Q2", "Q3"], "values": [100, 120, 130]}
#             })
#         elif content_type == 'story':
#             base_structure["slides"].append({
#                 "title": "Key Themes",
#                 "content": ["Main narrative themes", "Character development", "Plot highlights"],
#                 "type": "content",
#                 "chart_suggestion": "none",
#                 "chart_data": {}
#             })
#         elif content_type == 'sop':
#             base_structure["slides"].append({
#                 "title": "Procedure Overview",
#                 "content": ["Key procedures", "Compliance requirements", "Implementation steps"],
#                 "type": "content",
#                 "chart_suggestion": "table",
#                 "chart_data": {"headers": ["Step", "Description"], "rows": [["1", "Initiate process"], ["2", "Verify compliance"]]}
#             })
#         if tables:
#             base_structure["slides"].append({
#                 "title": "Data Summary",
#                 "content": [f"Table with {len(tables[0]['headers'])} columns and {len(tables[0]['rows'])} rows"],
#                 "type": "content",
#                 "chart_suggestion": "table",
#                 "chart_data": {"headers": tables[0]["headers"], "rows": tables[0]["rows"][:3]}
#             })
#         return base_structure

# class ChartGenerator:
#     def __init__(self):
#         self.colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
#         # Define base layouts for different styles
#         self.base_layout_configs = {
#             "default": {
#                 "font_size": 12,
#                 "title_font_size": 16,
#                 "bgcolor": 'white',
#                 "plot_bgcolor": 'white',
#                 "paper_bgcolor": 'white',
#                 "margin": dict(l=50, r=50, b=50, t=50),
#                 "legend_orientation": "h",
#                 "legend_xanchor": "center",
#                 "legend_x": 0.5,
#                 "legend_y": -0.2
#             },
#             "Minimalist": {
#                 "font_size": 10,
#                 "title_font_size": 14,
#                 "bgcolor": 'rgba(0,0,0,0)',
#                 "plot_bgcolor": 'rgba(0,0,0,0)',
#                 "paper_bgcolor": 'rgba(0,0,0,0)',
#                 "margin": dict(l=30, r=30, b=30, t=30),
#                 "showlegend": False,
#                 "xaxis_showgrid": False,
#                 "yaxis_showgrid": False
#             },
#             "Creative": {
#                 "font_size": 14,
#                 "title_font_size": 20,
#                 "bgcolor": '#f8f8f8',
#                 "plot_bgcolor": '#f8f8f8',
#                 "paper_bgcolor": '#f8f8f8',
#                 "margin": dict(l=60, r=60, b=60, t=60),
#                 "hovermode": "x unified"
#             }
#         }
    
#     def create_chart(self, chart_type: str, data: Dict[str, Any], style: str) -> str:
#         """Create chart using Plotly and return base64 encoded image"""
#         try:
#             fig = None
#             labels = data.get('labels', ['Category A', 'Category B', 'Category C']) 
#             values = data.get('values', [10, 20, 30])
#             headers = data.get('headers', [])
#             rows = data.get('rows', [])
            
#             # Apply base style configuration
#             layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])

#             plt.style.use('seaborn-v0_8')
#             if chart_type == 'bar':
#                 fig = px.bar(x=labels, y=values, color=labels, color_discrete_sequence=self.colors)
#                 fig.update_layout(title="<b>Data Visualization: Key Metrics</b>", **layout_config)
#             elif chart_type == 'line':
#                 fig = px.line(x=labels, y=values, markers=True, line_shape='spline', color_discrete_sequence=[self.colors[0]])
#                 fig.update_layout(title="<b>Trend Analysis Over Time</b>", **layout_config)
#             elif chart_type == 'pie':
#                 fig = px.pie(values=values, names=labels, color_discrete_sequence=self.colors, hole=0.3) 
#                 fig.update_layout(title="<b>Distribution Overview</b>", **layout_config)
#                 fig.update_traces(textinfo='percent+label', pull=[0.05 if i == values.index(max(values)) else 0 for i in range(len(values))]) 
#             elif chart_type == 'table':
#                 if headers and rows:
#                     fig = go.Figure(data=[go.Table(
#                         header=dict(values=headers,
#                                     fill_color=self.colors[0],
#                                     font=dict(color='white', size=layout_config["font_size"] + 2),
#                                     align='center',
#                                     height=30),
#                         cells=dict(values=[rows[i] for i in range(len(rows))],
#                                    fill_color='lavender',
#                                    font=dict(size=layout_config["font_size"]),
#                                    align='left',
#                                    height=25)
#                     )])
#                     fig.update_layout(title="<b>Detailed Data Table</b>", **layout_config)
            
#             if fig:
#                 buffer = io.BytesIO()
#                 fig.write_image(buffer, format='png', scale=2)
#                 buffer.seek(0)
#                 chart_base64 = base64.b64encode(buffer.read()).decode()
#                 return chart_base64
#             return ""
#         except Exception as e:
#             logger.error(f"Error creating chart: {str(e)}")
#             return ""

#     def create_statistical_chart(self, statistical_data: Dict[str, Any], chart_type: str, style: str) -> str:
#         """Create charts from extracted statistical data"""
#         try:
#             layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])
            
#             if chart_type == 'statistics_bar' and statistical_data['numbers']:
#                 # Create bar chart from numbers
#                 data = statistical_data['numbers'][:6]  # Limit to 6 items
#                 labels = [item['context'] for item in data]
#                 values = [item['value'] for item in data]
                
#                 fig = px.bar(x=labels, y=values, color=labels, 
#                             color_discrete_sequence=self.colors,
#                             title="<b>Key Metrics Overview</b>")
#                 fig.update_layout(**layout_config)
#                 fig.update_xaxis(tickangle=45)
                
#             elif chart_type == 'percentage_pie' and statistical_data['percentages']:
#                 # Create pie chart from percentages
#                 data = statistical_data['percentages'][:5]  # Limit to 5 items
#                 labels = [item['context'] for item in data]
#                 values = [item['value'] for item in data]
                
#                 fig = px.pie(values=values, names=labels, 
#                             color_discrete_sequence=self.colors,
#                             title="<b>Percentage Distribution</b>")
#                 fig.update_layout(**layout_config)
                
#             elif chart_type == 'monetary_line' and statistical_data['monetary_values']:
#                 # Create line chart from monetary values
#                 data = statistical_data['monetary_values'][:6]
#                 labels = [item['context'] for item in data]
#                 values = [item['value'] for item in data]
                
#                 fig = px.line(x=labels, y=values, markers=True,
#                             color_discrete_sequence=[self.colors[0]],
#                             title="<b>Financial Metrics Trend</b>")
#                 fig.update_layout(**layout_config)
#                 fig.update_xaxis(tickangle=45)
                
#             elif chart_type == 'comparison_bar' and statistical_data['comparisons']:
#                 # Create comparison chart
#                 data = statistical_data['comparisons'][:6]
#                 labels = [item['context'] for item in data]
#                 values = [item['value'] if item['type'] == 'increase' else -item['value'] for item in data]
#                 colors = ['green' if v > 0 else 'red' for v in values]
                
#                 fig = px.bar(x=labels, y=values, color=colors,
#                             title="<b>Performance Changes (%)</b>")
#                 fig.update_layout(**layout_config)
#                 fig.update_xaxis(tickangle=45)
                
#             else:
#                 return ""
                
#             # Convert to base64
#             buffer = io.BytesIO()
#             fig.write_image(buffer, format='png', scale=2)
#             buffer.seek(0)
#             return base64.b64encode(buffer.read()).decode()
            
#         except Exception as e:
#             logger.error(f"Error creating statistical chart: {str(e)}")
#             return ""

#     def suggest_chart_type(self, statistical_data: Dict[str, Any]) -> str:
#         """Suggest the best chart type based on available data"""
#         if statistical_data['percentages']:
#             return 'percentage_pie'
#         elif statistical_data['monetary_values']:
#             return 'monetary_line'
#         elif statistical_data['comparisons']:
#             return 'comparison_bar'
#         elif statistical_data['numbers']:
#             return 'statistics_bar'
#         return 'none'

# class PresentationGenerator:
#     def __init__(self, settings: Settings):
#         self.settings = settings
#         self.logger = logging.getLogger(__name__)
#         self.chart_generator = ChartGenerator()
#         self.default_theme = {
#             "colors": {
#                 "primary": "#003087",
#                 "secondary": "#FFFFFF",
#                 "accent": "#FFA500",
#                 "text": "#404040",
#                 "footer": "#808080"
#             },
#             "fonts": {
#                 "title": {"name": "Calibri", "size": 44, "bold": True},
#                 "subtitle": {"name": "Calibri", "size": 24, "bold": False},
#                 "content": {"name": "Calibri", "size": 18, "bold": False},
#                 "footer": {"name": "Calibri", "size": 10, "bold": False}
#             },
#             "footer": {
#                 "text": "{company_name} | {date} | Confidential",
#                 "position": {"left": 0.5, "top": 7.0, "width": 9.0, "height": 0.5}
#             }
#         }
#         self.theme = self.default_theme 
#         try:
#             theme_path = Path(__file__).parent / "templates" / "styles" / "corporate_theme.json"
#             if theme_path.exists():
#                 with open(theme_path, "r") as f:
#                     loaded_theme = json.load(f)
#                     self.theme = {**self.default_theme, **loaded_theme.get("theme", {})}
#                     for key in ["colors", "fonts", "footer"]:
#                         if key in loaded_theme.get("theme", {}) and isinstance(loaded_theme["theme"][key], dict):
#                             self.theme[key] = {**self.default_theme[key], **loaded_theme["theme"][key]}
#             else:
#                 self.logger.warning(f"Theme file not found at {theme_path}. Using default theme.")
#         except Exception as e:
#             self.logger.warning(f"Failed to load theme: {str(e)}. Using default theme.")
        
#         self.primary_color = RGBColor.from_string(self.theme["colors"]["primary"].lstrip("#"))
#         self.secondary_color = RGBColor.from_string(self.theme["colors"]["secondary"].lstrip("#"))
#         self.accent_color = RGBColor.from_string(self.theme["colors"]["accent"].lstrip("#"))
#         self.text_color = RGBColor.from_string(self.theme["colors"]["text"].lstrip("#"))
#         self.footer_color = RGBColor.from_string(self.theme["colors"]["footer"].lstrip("#"))
    

#     def _create_statistical_slide(self, prs: Presentation, statistical_data: Dict[str, Any], style: str):
#         """Create a slide specifically for statistical visualizations"""
#         slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
        
#         # Add title
#         title_placeholder = None
#         for shape in slide.placeholders:
#             if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:
#                 title_placeholder = shape
#                 break
        
#         if title_placeholder:
#             title_placeholder.text = "📊 Statistical Analysis"
#             title_placeholder.text_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["title"]["size"])
#             title_placeholder.text_frame.paragraphs[0].font.color.rgb = self.primary_color
#             title_placeholder.text_frame.paragraphs[0].font.bold = True
        
#         # Add statistical charts
#         chart_type = self.chart_generator.suggest_chart_type(statistical_data)
#         if chart_type != 'none':
#             chart_base64 = self.chart_generator.create_statistical_chart(statistical_data, chart_type, style)
#             if chart_base64:
#                 buffer = io.BytesIO(base64.b64decode(chart_base64))
#                 slide.shapes.add_picture(buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
        
#         # Add statistical summary text
#         summary_text = self._generate_statistical_summary(statistical_data)
#         if summary_text:
#             text_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
#             text_frame = text_box.text_frame
#             text_frame.text = summary_text
#             text_frame.paragraphs[0].font.size = Pt(14)
#             text_frame.paragraphs[0].font.color.rgb = self.text_color
        
#         self._add_footer(slide, "Statistical Analysis", f"{datetime.now().strftime('%B %Y')}")

#     def _generate_statistical_summary(self, statistical_data: Dict[str, Any]) -> str:
#         """Generate a text summary of statistical findings"""
#         summary_parts = []
        
#         if statistical_data['numbers']:
#             summary_parts.append(f"Key metrics: {len(statistical_data['numbers'])} data points identified")
        
#         if statistical_data['percentages']:
#             avg_percentage = sum(item['value'] for item in statistical_data['percentages']) / len(statistical_data['percentages'])
#             summary_parts.append(f"Average percentage value: {avg_percentage:.1f}%")
        
#         if statistical_data['monetary_values']:
#             total_value = sum(item['value'] for item in statistical_data['monetary_values'])
#             summary_parts.append(f"Total monetary value: ${total_value:,.0f}")
        
#         return " | ".join(summary_parts)

#     def generate_presentation(self, content: str, company_name: str, 
#                                 report_month: str, report_year: int, 
#                                 domain_focus: List[str], language: str,
#                                 content_type: str, tables: List[Dict[str, List]], style: str) -> str:
#         try:
#             try:
#                 prs = Presentation(str(self.settings.template_path))
#             except Exception as e:
#                 self.logger.warning(f"Template not found at {self.settings.template_path}: {str(e)}. Using default presentation.")
#                 st.warning(f"Using default presentation template as {self.settings.template_path} is missing.")
#                 prs = Presentation()
            
#             bedrock_client = BedrockClient(self.settings)
#             content_structure = bedrock_client.generate_presentation_content(
#                 content, company_name, report_month, report_year, domain_focus, language, content_type, tables
#             )
#             if not content_structure:
#                 content_structure = bedrock_client._create_default_structure(company_name, report_month, report_year, content_type, tables)
            
#             self._create_enhanced_title_slide(prs, content_structure['title'], company_name, report_month, report_year, domain_focus, style)
#             for slide_data in content_structure['slides']:
#                 self._create_enhanced_content_slide(prs, slide_data, style)
            
#             filename = f"Presentation_{company_name.replace(' ', '_')}_{report_month}_{report_year}.pptx"
#             output_path = self.settings.output_folder / filename
#             prs.save(str(output_path))
#             return str(output_path)
#         except Exception as e:
#             self.logger.error(f"Error generating presentation: {str(e)}")
#             st.error(f"Error generating presentation: {str(e)}")
#             raise
    
#     def _create_enhanced_title_slide(self, prs: Presentation, title: str, company_name: str, 
#                                       report_month: str, report_year: int, domain_focus: List[str], style: str):
#         slide_layout = prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)

#         title_placeholder = None
#         subtitle_placeholder = None

#         for shape in slide.placeholders:
#             if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title placeholder
#                 title_placeholder = shape
#             elif hasattr(shape, 'name') and shape.name == "Subtitle 2":
#                 subtitle_placeholder = shape
        
#         # Populate title
#         if title_placeholder:
#             text_frame = title_placeholder.text_frame
#             text_frame.text = title
#             paragraph = text_frame.paragraphs[0]
#             paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
#             paragraph.font.color.rgb = self.primary_color
#             paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
#             paragraph.alignment = PP_ALIGN.CENTER
#         else:
#             left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(2)
#             title_box = slide.shapes.add_textbox(left, top, width, height)
#             title_frame = title_box.text_frame
#             title_frame.text = title
#             title_paragraph = title_frame.paragraphs[0]
#             title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
#             title_paragraph.font.color.rgb = self.primary_color
#             title_paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
#             title_paragraph.alignment = PP_ALIGN.CENTER
        
#         # Populate subtitle
#         subtitle_text = f"{company_name} | {report_month} {report_year} | {', '.join(domain_focus)}"
#         if subtitle_placeholder:
#             text_frame = subtitle_placeholder.text_frame
#             text_frame.text = subtitle_text
#             paragraph = text_frame.paragraphs[0]
#             paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
#             paragraph.font.color.rgb = self.accent_color
#             paragraph.alignment = PP_ALIGN.CENTER
#         else:
#             subtitle_top = Inches(3.5)
#             subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1.5))
#             subtitle_frame = subtitle_box.text_frame
#             subtitle_frame.text = subtitle_text
#             subtitle_paragraph = subtitle_frame.paragraphs[0]
#             subtitle_paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
#             subtitle_paragraph.font.color.rgb = self.accent_color
#             subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
#         if style in ["Creative", "Minimalist"]:
#             self._add_decorative_image(slide, style)
#         self._add_footer(slide, company_name, f"{report_month} {report_year}")
    
#     def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], style: str):
#         slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
#         slide = prs.slides.add_slide(slide_layout)
        
#         title_placeholder = None
#         body_placeholder = None
#         for shape in slide.placeholders:
#             if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title placeholder
#                 title_placeholder = shape
#             elif hasattr(shape, 'has_text_frame') and shape.has_text_frame and hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY placeholder
#                 body_placeholder = shape

#         # Populate title
#         if title_placeholder:
#             title_shape = title_placeholder
#             title_shape.text = slide_data['title']
#             title_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["title"]["size"])
#             title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
#             title_shape.text_frame.paragraphs[0].font.bold = True
#         else:
#             left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
#             title_box = slide.shapes.add_textbox(left, top, width, height)
#             title_frame = title_box.text_frame
#             title_frame.text = slide_data['title']
#             title_paragraph = title_frame.paragraphs[0]
#             title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
#             title_paragraph.font.color.rgb = self.primary_color
#             title_paragraph.font.bold = True
#             title_paragraph.alignment = PP_ALIGN.LEFT
        
#         # Populate content (bullet points or chart)
#         if slide_data.get('chart_suggestion') != "none" and slide_data.get('chart_data'):
#             content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
#             content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#             self._add_bullet_points(content_box.text_frame, slide_data['content'])

#             chart_left = Inches(5.0)
#             self._add_chart_to_slide(slide, slide_data['chart_suggestion'], slide_data['chart_data'], chart_left, content_top, style)
#         else:
#             if body_placeholder:
#                 content_box = body_placeholder
#                 # Clear existing paragraphs in the placeholder's text frame
#                 for para in list(content_box.text_frame.paragraphs):
#                     content_box.text_frame.remove_paragraph(para)
#                 self._add_bullet_points(content_box.text_frame, slide_data['content'])
#             else:
#                 content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
#                 content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
#                 self._add_bullet_points(content_box.text_frame, slide_data['content'])
        
#         self._add_footer(slide, "Your Company", f"{datetime.now().strftime('%B %Y')}")
    
#     def _add_bullet_points(self, text_frame, content_list: List[str]):
#     # Clear existing content
#         text_frame.clear()
        
#         # Add content
#         for i, item in enumerate(content_list):
#             if i == 0:
#                 # Use the first paragraph that already exists
#                 paragraph = text_frame.paragraphs[0]
#             else:
#                 # Add new paragraphs for additional items
#                 paragraph = text_frame.add_paragraph()
            
#             paragraph.text = item
#             paragraph.font.size = Pt(self.theme["fonts"]["content"]["size"])
#             paragraph.font.color.rgb = self.text_color
#             paragraph.level = 0
#             paragraph.space_after = Pt(12)


#     def _add_chart_to_slide(self, slide, chart_type: str, chart_data: Dict[str, Any], left: float, top: float, style: str):
#         try:
#             chart_width, chart_height = Inches(4.5), Inches(3.5)
#             if chart_type == 'table':
#                 table = slide.shapes.add_table(
#                     len(chart_data.get('rows', [])[:5]) + 1, # Max 5 rows for simplicity + header
#                     len(chart_data.get('headers', [])),
#                     left, top, chart_width, chart_height
#                 ).table
#                 for i, header in enumerate(chart_data.get('headers', [])):
#                     table.cell(0, i).text = header
#                     table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)
#                 for i, row in enumerate(chart_data.get('rows', [])[:5], 1): # Limit rows to avoid overflow
#                     for j, value in enumerate(row):
#                         table.cell(i, j).text = str(value)
#                         table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
#             else:
#                 chart_base64 = self.chart_generator.create_chart(chart_type, chart_data, style)
#                 if chart_base64:
#                     buffer = io.BytesIO(base64.b64decode(chart_base64))
#                     slide.shapes.add_picture(buffer, left, top, chart_width, chart_height)
#         except Exception as e:
#             self.logger.error(f"Error adding chart: {str(e)}")
#             st.error(f"Error adding chart: {str(e)}")
#             fallback_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
#             fallback_box.text_frame.text = "📊 Chart unavailable\n(Data visualization)"
#             fallback_box.text_frame.paragraphs[0].font.size = Pt(16)
#             fallback_box.text_frame.paragraphs[0].font.color.rgb = self.accent_color
    
#     def _add_decorative_image(self, slide, style: str):
#         """Add a decorative image or icon based on style"""
#         try:
#             img_filename = "creative_icon.png" if style == "Creative" else "minimal_icon.png"
#             img_path = Path(__file__).parent / "templates" / "icons" / img_filename
            
#             if not img_path.exists():
#                 img = Image.new('RGB', (200, 200), color=(self.accent_color.rgb[0], self.accent_color.rgb[1], self.accent_color.rgb[2]))
#                 img.save(img_path)

#             left, top = Inches(8.5), Inches(0.5)
#             width = Inches(1.5) if style == "Creative" else Inches(1)
#             height = Inches(1.5) if style == "Creative" else Inches(1)

#             slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)

#         except Exception as e:
#             self.logger.error(f"Error adding decorative image: {str(e)}")
#             left, top, width, height = Inches(8.5), Inches(0.5), Inches(1), Inches(1)
#             shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
#             shape.fill.solid()
#             shape.fill.fore_color.rgb = self.accent_color
#             shape.line.fill.background()
    
#     def _add_footer(self, slide, company_name: str, date: str):
#         left = Inches(self.theme["footer"]["position"]["left"])
#         top = Inches(self.theme["footer"]["position"]["top"])
#         width = Inches(self.theme["footer"]["position"]["width"])
#         height = Inches(self.theme["footer"]["position"]["height"])

#         # Add shape for footer background
#         footer_background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(10), height)
#         footer_background.fill.solid()
#         footer_background.fill.fore_color.rgb = self.secondary_color
#         footer_background.line.fill.background()

#         footer_box = slide.shapes.add_textbox(left, top, width, height)
#         footer_frame = footer_box.text_frame
#         footer_frame.text = self.theme["footer"]["text"].format(company_name=company_name, date=date)
#         footer_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
#         footer_frame.paragraphs[0].font.color.rgb = self.footer_color
#         footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

#         # Add slide number
#         slide_number_left = Inches(0.5)
#         slide_number_box = slide.shapes.add_textbox(slide_number_left, top, Inches(1), height)
#         slide_number_frame = slide_number_box.text_frame
#         # You might want to implement actual slide numbering if you prefer
#         # This uses PowerPoint's internal slide ID which is not sequential
#         slide_number_frame.text = f"Page {slide.slide_id}" 
#         slide_number_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
#         slide_number_frame.paragraphs[0].font.color.rgb = self.footer_color
#         slide_number_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# @st.cache_resource
# def initialize_components():
#     settings = Settings()
#     doc_processor = DocumentProcessor()
#     presentation_generator = PresentationGenerator(settings)
#     return settings, doc_processor, presentation_generator

# def main():
#     st.markdown("""
#     <div class="main-header">
#         <h1>🎯 Smart Presentation Generator</h1>
#         <p>Transform any content into beautiful, professional presentations with AI-powered insights and visualizations</p>
#     </div>
#     """, unsafe_allow_html=True)
    
#     settings, doc_processor, presentation_generator = initialize_components()
#     content_type = "Not yet processed"
    
#     with st.sidebar:
#         st.header("⚙️ Configuration")
#         input_method = st.radio(
#             "Choose input method:",
#             ["📄 Upload PDF", "📝 Direct Text Input"],
#             horizontal=True
#         )
#         target_language = st.selectbox(
#             "Target Language",
#             options=["en", "es", "fr", "de", "it", "pt", "nl", "ru", "zh", "ja", "ko", "ar", "hi"],
#             format_func=lambda x: {
#                 "en": "🇺🇸 English", "es": "🇪🇸 Spanish", "fr": "🇫🇷 French", 
#                 "de": "🇩🇪 German", "it": "🇮🇹 Italian", "pt": "🇵🇹 Portuguese", 
#                 "nl": "🇳🇱 Dutch", "ru": "🇷🇺 Russian", "zh": "🇨🇳 Chinese", 
#                 "ja": "🇯🇵 Japanese", "ko": "🇰🇷 Korean", "ar": "🇸🇦 Arabic", 
#                 "hi": "🇮🇳 Hindi"
#             }[x]
#         )
#         company_name = st.text_input("Company Name", "Your Company")
#         col1, col2 = st.columns(2)
#         with col1:
#             report_month = st.selectbox(
#                 "Month",
#                 options=["January", "February", "March", "April", "May", "June",
#                          "July", "August", "September", "October", "November", "December"],
#                 index=datetime.now().month - 1
#             )
#         with col2:
#             report_year = st.number_input("Year", min_value=2020, max_value=2030, value=datetime.now().year)
#         domain_focus = st.multiselect(
#             "Focus Areas",
#             options=["Logistics Vehicles", "Cargo Management", "Fleet Operations", 
#                      "Supply Chain", "Transportation Efficiency", "Business Performance",
#                      "Financial Reporting", "Market Analysis", "Strategic Planning",
#                      "Narrative Development", "Character Analysis", "Plot Structure",
#                      "Standard Operating Procedures", "Compliance and Regulations", "Operational Guidelines"],
#             default=["Business Performance"]
#         )
#         presentation_style = st.selectbox(
#             "Presentation Style",
#             options=["Default", "Minimalist", "Creative"]
#         )
        
#         st.markdown("---")
#         st.markdown("Developed by **Your Team Name**")

#     document_content = ""
#     tables_extracted: List[Dict[str, List]] = []

#     if input_method == "📄 Upload PDF":
#         uploaded_file = st.file_uploader("Upload a PDF document", type=["pdf"])
#         if uploaded_file:
#             with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
#                 tmp_file.write(uploaded_file.getvalue())
#                 temp_pdf_path = tmp_file.name
            
#             with st.spinner("Extracting text and analyzing content..."):
#                 document_content = doc_processor.extract_text_from_pdf(temp_pdf_path)
#                 tables_extracted = doc_processor.extract_tables(document_content)
#                 content_analysis = doc_processor.analyze_content_type(document_content)
#                 content_type = content_analysis['type']
#                 st.session_state['document_content'] = document_content
#                 st.session_state['tables_extracted'] = tables_extracted
#                 st.session_state['content_type'] = content_type
            
#             os.unlink(temp_pdf_path) # Clean up temp file
#             if document_content:
#                 st.success("PDF processed successfully!")
#                 st.expander("Extracted Text Preview").write(document_content[:1000] + "...")
#                 if tables_extracted:
#                     st.expander("Extracted Tables Preview").write(tables_extracted)
#             else:
#                 st.error("Failed to extract content from PDF.")

#     elif input_method == "📝 Direct Text Input":
#         text_input = st.text_area("Paste your content here:", height=300)
#         if text_input:
#             with st.spinner("Analyzing content..."):
#                 document_content = doc_processor.process_text_document(text_input)
#                 tables_extracted = doc_processor.extract_tables(document_content)
#                 content_analysis = doc_processor.analyze_content_type(document_content)
#                 content_type = content_analysis['type']
#                 st.session_state['document_content'] = document_content
#                 st.session_state['tables_extracted'] = tables_extracted
#                 st.session_state['content_type'] = content_type
#             st.success("Text content processed successfully!")

#     if st.button("Generate Presentation 🚀", type="primary"):
#         if 'document_content' not in st.session_state or not st.session_state['document_content']:
#             st.error("Please upload a document or enter text content first.")
#             return

#         with st.spinner("Generating professional presentation... This may take a few moments."):
#             try:
#                 # Retrieve processed content from session state
#                 doc_content = st.session_state['document_content']
#                 tables = st.session_state['tables_extracted']
#                 c_type = st.session_state['content_type']

#                 # Translate content if target language is not English
#                 if target_language != 'en':
#                     bedrock_client = BedrockClient(settings)
#                     doc_content = bedrock_client.translate_text(doc_content, target_language)
#                     # Note: Translating tables directly via Bedrock might be complex
#                     # For simplicity, tables remain in original language or require separate translation logic
                
#                 output_ppt_path = presentation_generator.generate_presentation(
#                     doc_content, company_name, report_month, report_year, domain_focus, target_language,
#                     c_type, tables, presentation_style
#                 )
                
#                 if output_ppt_path and os.path.exists(output_ppt_path):
#                     with open(output_ppt_path, "rb") as file:
#                         btn = st.download_button(
#                             label="Download Presentation ⬇️",
#                             data=file,
#                             file_name=os.path.basename(output_ppt_path),
#                             mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
#                         )
#                     st.success("Presentation generated successfully! Click the button above to download.")
#                     st.balloons()
#                 else:
#                     st.error("Failed to generate presentation. Please check logs for details.")
#             except Exception as e:
#                 st.error(f"An unexpected error occurred during presentation generation: {e}")
#                 logger.error(f"Critical error in main generation flow: {e}")

#     st.markdown("---")
#     st.info("💡 **Tips for best results:**\n"
#             "- Upload clear, well-formatted PDF documents.\n"
#             "- Specify relevant focus areas for tailored content.\n"
#             "- Ensure your AWS Bedrock credentials are set up in `.env` if you're using LLM functionality.")

# if __name__ == "__main__":
#     main()


import streamlit as st
import os
import tempfile
import json
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional, Union
import zipfile
import io
import base64
import re
import boto3
from botocore.exceptions import ClientError
import PyPDF2
from langdetect import detect
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from dotenv import load_dotenv
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import numpy as np
from PIL import Image

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(), logging.FileHandler("app.log")]
)
logger = logging.getLogger(__name__)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        padding: 20px;
        border-radius: 10px;
        color: white;
        margin-bottom: 30px;
        text-align: center;
    }
    .feature-box {
        background: #f8f9fa;
        padding: 15px;
        border-radius: 8px;
        border-left: 4px solid #667eea;
        margin: 10px 0;
    }
    .metric-container {
        background: white;
        padding: 20px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        margin: 10px 0;
    }
    .preview-image {
        max-width: 100%;
        border-radius: 8px;
        margin-top: 10px;
    }
</style>
""", unsafe_allow_html=True)

class Settings:
    def __init__(self):
        load_dotenv()
        self.aws_access_key_id = os.getenv('AWS_ACCESS_KEY_ID')
        self.aws_secret_access_key = os.getenv('AWS_SECRET_ACCESS_KEY')
        self.aws_session_token = os.getenv('AWS_SESSION_TOKEN')
        self.aws_default_region = os.getenv('AWS_DEFAULT_REGION', 'us-east-1')
        self.bedrock_model_id = os.getenv('BEDROCK_MODEL_ID', 'anthropic.claude-3-sonnet-20240229-v1:0')
        self.bedrock_region = os.getenv('BEDROCK_REGION', 'us-east-1')
        
        # Robust template path resolution
        self.template_path = Path(__file__).parent / 'templates' / 'monthly_template.pptx'
        
        self.upload_folder = Path('uploads')
        self.output_folder = Path('output')
        self.max_file_size = 10485760  # 10MB
        self.allowed_extensions = ['pdf']
        self.default_source_language = 'auto'
        self.default_target_language = 'en'
        self.supported_languages = ['en', 'es', 'fr', 'de', 'it', 'pt', 'nl', 'ru', 'zh', 'ja', 'ko', 'ar', 'hi']
        
        # Create directories, including templates/styles if they don't exist
        self.upload_folder.mkdir(exist_ok=True)
        self.output_folder.mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates').mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates' / 'styles').mkdir(exist_ok=True)
        (Path(__file__).parent / 'templates' / 'icons').mkdir(exist_ok=True) # Ensure icons folder exists

        # Set up AWS session
        try:
            boto3.setup_default_session(
                aws_access_key_id=self.aws_access_key_id,
                aws_secret_access_key=self.aws_secret_access_key,
                aws_session_token=self.aws_session_token,
                region_name=self.aws_default_region
            )
        except Exception as e:
            st.warning(f"AWS setup warning: {str(e)}")
            logger.error(f"AWS setup error: {str(e)}")

class DocumentProcessor:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        try:
            text = ""
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
            return self._clean_text(text)
        except Exception as e:
            self.logger.error(f"Error extracting text from PDF: {str(e)}")
            st.error(f"Error extracting PDF: {str(e)}")
            return ""
    
    def process_text_document(self, content: str) -> str:
        return self._clean_text(content)
    
    def _clean_text(self, text: str) -> str:
        text = ' '.join(text.split())
        text = text.replace('\x00', '').replace('\ufeff', '')
        return text
    
    def detect_language(self, text: str) -> Optional[str]:
        try:
            if len(text.strip()) < 10:
                return None
            return detect(text)
        except:
            return None
    
    def extract_tables(self, text: str) -> List[Dict[str, List]]:
        """Extract table-like data from text"""
        tables = []
        lines = text.split('\n')
        current_table = None
        headers = None
        
        for line in lines:
            if re.search(r'[,|]\s*|\d+\s+\w+\s+\d+', line):
                items = [item.strip() for item in re.split(r'[,|]\s*|\s{2,}', line) if item.strip()]
                if len(items) >= 2:
                    if not headers:
                        headers = items
                        current_table = {"headers": headers, "rows": []}
                    else:
                        if len(items) == len(headers):
                            current_table["rows"].append(items)
            else:
                if current_table:
                    tables.append(current_table)
                    headers = None
                    current_table = None
        
        if current_table:
            tables.append(current_table)
        
        return tables
    
    def extract_statistical_data(self, text: str) -> Dict[str, Any]:
        """Extract numerical data, percentages, and statistical information from text"""
        statistical_data = {
            'numbers': [],
            'percentages': [],
            'dates': [],
            'monetary_values': [],
            'metrics': [],
            'comparisons': []
        }
        
        # Extract numbers with context
        number_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:,\d{3})*(?:\.\d+)?)\s*(%|percent|million|billion|thousand|units|kg|tons|miles|km)?'
        numbers = re.findall(number_pattern, text, re.IGNORECASE)
        
        for context, value, unit in numbers:
            try:
                numeric_value = float(value.replace(',', ''))
                statistical_data['numbers'].append({
                    'context': context.strip(),
                    'value': numeric_value,
                    'unit': unit,
                    'original': f"{value} {unit}".strip()
                })
            except ValueError:
                continue
        
        # Extract percentages
        percentage_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*(\d+(?:\.\d+)?)\s*(%|percent)'
        percentages = re.findall(percentage_pattern, text, re.IGNORECASE)
        
        for context, value, unit in percentages:
            try:
                statistical_data['percentages'].append({
                    'context': context.strip(),
                    'value': float(value),
                    'original': f"{value}%"
                })
            except ValueError:
                continue
        
        # Extract monetary values
        money_pattern = r'(\w+(?:\s+\w+){0,2})\s*:?\s*\$(\d+(?:,\d{3})*(?:\.\d+)?)\s*(million|billion|thousand|M|B|K)?'
        monetary = re.findall(money_pattern, text, re.IGNORECASE)
        
        for context, value, multiplier in monetary:
            try:
                numeric_value = float(value.replace(',', ''))
                multiplier_map = {'million': 1000000, 'M': 1000000, 'billion': 1000000000, 'B': 1000000000, 'thousand': 1000, 'K': 1000}
                if multiplier and multiplier in multiplier_map:
                    numeric_value *= multiplier_map[multiplier]
                
                statistical_data['monetary_values'].append({
                    'context': context.strip(),
                    'value': numeric_value,
                    'original': f"${value} {multiplier}".strip()
                })
            except ValueError:
                continue
        
        # Extract comparison data (increased by X%, decreased by Y%)
        comparison_pattern = r'(\w+(?:\s+\w+){0,2})\s*(?:increased|decreased|rose|fell|grew|declined)\s*(?:by\s*)?(\d+(?:\.\d+)?)\s*(%|percent)'
        comparisons = re.findall(comparison_pattern, text, re.IGNORECASE)
        
        for context, value, unit in comparisons:
            try:
                statistical_data['comparisons'].append({
                    'context': context.strip(),
                    'value': float(value),
                    'type': 'increase' if any(word in text.lower() for word in ['increased', 'rose', 'grew']) else 'decrease'
                })
            except ValueError:
                continue
        
        return statistical_data
    
    def analyze_content_type(self, text: str) -> Dict[str, Any]:
        content_analysis = {
            'type': 'general',
            'has_numbers': bool(re.search(r'\d+', text)),
            'has_dates': bool(re.search(r'\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b', text, re.I)),
            'has_tables': bool(self.extract_tables(text)),
            'topics': [],
            'potential_charts': []
        }
        text_lower = text.lower()
        logistics_keywords = ['logistics', 'cargo', 'fleet', 'vehicle', 'transport', 'shipping', 'warehouse']
        business_keywords = ['revenue', 'profit', 'sales', 'market', 'customer', 'performance']
        story_keywords = ['story', 'narrative', 'character', 'plot', 'once upon']
        sop_keywords = ['procedure', 'sop', 'standard operating', 'guideline', 'protocol', 'instruction']
        
        if any(keyword in text_lower for keyword in logistics_keywords):
            content_analysis['type'] = 'logistics'
            content_analysis['topics'] = ['Fleet Operations', 'Cargo Management', 'Delivery Metrics']
        elif any(keyword in text_lower for keyword in business_keywords):
            content_analysis['type'] = 'business'
            content_analysis['topics'] = ['Financials', 'Market Trends', 'Performance']
        elif any(keyword in text_lower for keyword in story_keywords):
            content_analysis['type'] = 'story'
            content_analysis['topics'] = ['Narrative', 'Themes', 'Characters']
        elif any(keyword in text_lower for keyword in sop_keywords):
            content_analysis['type'] = 'sop'
            content_analysis['topics'] = ['Procedures', 'Guidelines', 'Compliance']
        
        if content_analysis['has_numbers'] or content_analysis['has_tables']:
            content_analysis['potential_charts'] = ['bar', 'line', 'pie', 'table']
        
        # Add statistical analysis
        statistical_data = self.extract_statistical_data(text)
        content_analysis['statistical_data'] = statistical_data
        content_analysis['has_statistics'] = len(statistical_data['numbers']) > 0 or len(statistical_data['percentages']) > 0

        # Suggest chart types based on data
        if statistical_data['percentages']:
            content_analysis['potential_charts'].extend(['pie', 'donut'])
        if statistical_data['comparisons']:
            content_analysis['potential_charts'].extend(['bar', 'column'])
        if statistical_data['monetary_values']:
            content_analysis['potential_charts'].extend(['line', 'area'])
        
        return content_analysis

class BedrockClient:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.logger = logging.getLogger(__name__)
        try:
            self.bedrock_client = boto3.client('bedrock-runtime', region_name=self.settings.bedrock_region)
        except Exception as e:
            self.logger.error(f"Failed to initialize Bedrock client: {str(e)}")
            st.error(f"Failed to initialize Bedrock client: {str(e)}")
            self.bedrock_client = None
    
    def invoke_model(self, prompt: str, max_tokens: int = 4000) -> Optional[str]:
        if not self.bedrock_client:
            st.warning("Bedrock client unavailable. Using fallback response.")
            return self._fallback_response(prompt)
        
        try:
            body = {
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": max_tokens,
                "messages": [{"role": "user", "content": prompt}]
            }
            response = self.bedrock_client.invoke_model(
                modelId=self.settings.bedrock_model_id,
                body=json.dumps(body)
            )
            response_body = json.loads(response['body'].read())
            return response_body.get('content', [{}])[0].get('text', '')
        except ClientError as e:
            self.logger.error(f"Bedrock client error: {str(e)}")
            st.error(f"AWS Bedrock Error: {str(e)}")
            return self._fallback_response(prompt)
        except Exception as e:
            self.logger.error(f"Error invoking Bedrock model: {str(e)}")
            st.error(f"Error invoking Bedrock: {str(e)}")
            return self._fallback_response(prompt)
    
    def _fallback_response(self, prompt: str) -> str:
        if "translate" in prompt.lower():
            return "Translation not available - using original text"
        elif "json" in prompt.lower():
            return json.dumps({
                "title": "Generated Presentation",
                "slides": [
                    {
                        "title": "Overview",
                        "content": ["Key insights from the document", "Main topics covered", "Summary of findings"],
                        "type": "content",
                        "chart_suggestion": "none",
                        "chart_data": {}
                    }
                ]
            })
        return "Content analysis not available - using default structure"
    
    def translate_text(self, text: str, target_language: str, source_language: str = 'auto') -> str:
        if source_language == target_language or not text:
            return text
        prompt = f"""
        Translate the following text to {target_language}, maintaining structure and formatting:
        {text[:4000]}
        Translation:
        """
        response = self.invoke_model(prompt)
        return response or text
    
    def generate_presentation_content(self, text: str, company_name: str, 
                                     report_month: str, report_year: int, 
                                     domain_focus: List[str], language: str,
                                     content_type: str, tables: List[Dict[str, List]], 
                                     content_analysis: Dict[str, Any] = None) -> Dict[str, Any]:
        focus_areas = ", ".join(domain_focus)
        table_summary = ""
        if tables:
            table_summary = "\nTables found:\n"
            for i, table in enumerate(tables, 1):
                table_summary += f"Table {i}:\nHeaders: {', '.join(table['headers'])}\nRows: {len(table['rows'])}\n"
        
        prompt = f"""
        Analyze this {content_type} document and create a professional presentation structure for {company_name} 
        for {report_month} {report_year} in {language}. Focus areas: {focus_areas}.
        
        Content: {text[:4000]}
        {table_summary}
        
        Steps:
        1. Identify key topics, metrics, and insights from the text and tables.
        2. Group information into 5-7 logical slides, including a title slide.
        3. For each slide, provide a title, 3-5 bullet points, and suggest a chart type (bar, line, pie, table, or none) with sample data.
        4. Tailor content to the {content_type} context, using tables for data-driven slides if available.
        
        Return JSON format:
        {{
            "title": "Presentation title",
            "slides": [
                {{
                    "title": "Slide title",
                    "content": ["point 1", "point 2", "point 3"],
                    "type": "content",
                    "chart_suggestion": "bar|line|pie|table|none",
                    "chart_data": {{"labels": ["A", "B"], "values": [10, 20]}} or {{"headers": [], "rows": [[]]}}
                }}
            ]
        }}
        """
        # Add statistical data to prompt if available
        if content_analysis and content_analysis.get('has_statistics', False):
            statistical_data = content_analysis['statistical_data']
            stats_summary = f"""
            Statistical data found:
            - Numbers: {len(statistical_data['numbers'])} metrics
            - Percentages: {len(statistical_data['percentages'])} values
            - Monetary values: {len(statistical_data['monetary_values'])} amounts
            - Comparisons: {len(statistical_data['comparisons'])} changes
            """
            prompt += stats_summary
        
        response = self.invoke_model(prompt, max_tokens=6000)
        if response:
            try:
                json_start = response.find('{')
                json_end = response.rfind('}') + 1
                json_str = response[json_start:json_end]
                return json.loads(json_str)
            except Exception as e:
                self.logger.error(f"Error parsing JSON: {str(e)}")
                st.error(f"Error parsing Bedrock response: {str(e)}")
        return self._create_default_structure(company_name, report_month, report_year, content_type, tables)
    
    def _create_default_structure(self, company_name: str, report_month: str, 
                                 report_year: int, content_type: str, tables: List[Dict[str, List]]) -> Dict[str, Any]:
        base_structure = {
            "title": f"{company_name} - {report_month} {report_year} Report",
            "slides": [
                {
                    "title": "Executive Summary",
                    "content": ["Overview of key findings", "Main topics identified", "Summary of performance"],
                    "type": "content",
                    "chart_suggestion": "none",
                    "chart_data": {}
                },
                {
                    "title": "Key Insights",
                    "content": ["Primary observations", "Notable trends", "Actionable takeaways"],
                    "type": "content",
                    "chart_suggestion": "none",
                    "chart_data": {}
                }
            ]
        }
        if content_type == 'logistics':
            base_structure["slides"].append({
                "title": "Logistics Metrics",
                "content": ["Operational efficiency", "Fleet performance", "Delivery metrics"],
                "type": "content",
                "chart_suggestion": "bar",
                "chart_data": {"labels": ["Efficiency", "On-Time", "Cost"], "values": [85, 92, 75]}
            })
        elif content_type == 'business':
            base_structure["slides"].append({
                "title": "Financial Overview",
                "content": ["Revenue trends", "Profit margins", "Cost analysis"],
                "type": "content",
                "chart_suggestion": "line",
                "chart_data": {"labels": ["Q1", "Q2", "Q3"], "values": [100, 120, 130]}
            })
        elif content_type == 'story':
            base_structure["slides"].append({
                "title": "Key Themes",
                "content": ["Main narrative themes", "Character development", "Plot highlights"],
                "type": "content",
                "chart_suggestion": "none",
                "chart_data": {}
            })
        elif content_type == 'sop':
            base_structure["slides"].append({
                "title": "Procedure Overview",
                "content": ["Key procedures", "Compliance requirements", "Implementation steps"],
                "type": "content",
                "chart_suggestion": "table",
                "chart_data": {"headers": ["Step", "Description"], "rows": [["1", "Initiate process"], ["2", "Verify compliance"]]}
            })
        if tables:
            base_structure["slides"].append({
                "title": "Data Summary",
                "content": [f"Table with {len(tables[0]['headers'])} columns and {len(tables[0]['rows'])} rows"],
                "type": "content",
                "chart_suggestion": "table",
                "chart_data": {"headers": tables[0]["headers"], "rows": tables[0]["rows"][:3]}
            })
        return base_structure

class ChartGenerator:
    def __init__(self):
        self.colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
        # Define base layouts for different styles
        self.base_layout_configs = {
            "default": {
                "font_size": 12,
                "title_font_size": 16,
                "bgcolor": 'white',
                "plot_bgcolor": 'white',
                "paper_bgcolor": 'white',
                "margin": dict(l=50, r=50, b=50, t=50),
                "legend_orientation": "h",
                "legend_xanchor": "center",
                "legend_x": 0.5,
                "legend_y": -0.2
            },
            "Minimalist": {
                "font_size": 10,
                "title_font_size": 14,
                "bgcolor": 'rgba(0,0,0,0)',
                "plot_bgcolor": 'rgba(0,0,0,0)',
                "paper_bgcolor": 'rgba(0,0,0,0)',
                "margin": dict(l=30, r=30, b=30, t=30),
                "showlegend": False,
                "xaxis_showgrid": False,
                "yaxis_showgrid": False
            },
            "Creative": {
                "font_size": 14,
                "title_font_size": 20,
                "bgcolor": '#f8f8f8',
                "plot_bgcolor": '#f8f8f8',
                "paper_bgcolor": '#f8f8f8',
                "margin": dict(l=60, r=60, b=60, t=60),
                "hovermode": "x unified"
            }
        }
    
    def create_chart(self, chart_type: str, data: Dict[str, Any], style: str) -> str:
        """Create chart using Plotly and return base64 encoded image"""
        try:
            fig = None
            labels = data.get('labels', ['Category A', 'Category B', 'Category C']) 
            values = data.get('values', [10, 20, 30])
            headers = data.get('headers', [])
            rows = data.get('rows', [])
            
            # Apply base style configuration
            layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])

            plt.style.use('seaborn-v0_8')
            if chart_type == 'bar':
                fig = px.bar(x=labels, y=values, color=labels, color_discrete_sequence=self.colors)
                fig.update_layout(title="<b>Data Visualization: Key Metrics</b>", **layout_config)
            elif chart_type == 'line':
                fig = px.line(x=labels, y=values, markers=True, line_shape='spline', color_discrete_sequence=[self.colors[0]])
                fig.update_layout(title="<b>Trend Analysis Over Time</b>", **layout_config)
            elif chart_type == 'pie':
                fig = px.pie(values=values, names=labels, color_discrete_sequence=self.colors, hole=0.3) 
                fig.update_layout(title="<b>Distribution Overview</b>", **layout_config)
                fig.update_traces(textinfo='percent+label', pull=[0.05 if i == values.index(max(values)) else 0 for i in range(len(values))]) 
            elif chart_type == 'table':
                if headers and rows:
                    fig = go.Figure(data=[go.Table(
                        header=dict(values=headers,
                                    fill_color=self.colors[0],
                                    font=dict(color='white', size=layout_config["font_size"] + 2),
                                    align='center',
                                    height=30),
                        cells=dict(values=[rows[i] for i in range(len(rows))],
                                   fill_color='lavender',
                                   font=dict(size=layout_config["font_size"]),
                                   align='left',
                                   height=25)
                    )])
                    fig.update_layout(title="<b>Detailed Data Table</b>", **layout_config)
            
            if fig:
                buffer = io.BytesIO()
                fig.write_image(buffer, format='png', scale=2)
                buffer.seek(0)
                chart_base64 = base64.b64encode(buffer.read()).decode()
                return chart_base64
            return ""
        except Exception as e:
            logger.error(f"Error creating chart: {str(e)}")
            return ""
    
    def create_statistical_chart(self, statistical_data: Dict[str, Any], chart_type: str, style: str) -> str:
        """Create charts from extracted statistical data"""
        try:
            layout_config = self.base_layout_configs.get(style, self.base_layout_configs["default"])
            
            if chart_type == 'statistics_bar' and statistical_data['numbers']:
                # Create bar chart from numbers
                data = statistical_data['numbers'][:6]  # Limit to 6 items
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
                
                fig = px.bar(x=labels, y=values, color=labels, 
                            color_discrete_sequence=self.colors,
                            title="<b>Key Metrics Overview</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
                
            elif chart_type == 'percentage_pie' and statistical_data['percentages']:
                # Create pie chart from percentages
                data = statistical_data['percentages'][:5]  # Limit to 5 items
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
                
                fig = px.pie(values=values, names=labels, 
                            color_discrete_sequence=self.colors,
                            title="<b>Percentage Distribution</b>")
                fig.update_layout(**layout_config)
                
            elif chart_type == 'monetary_line' and statistical_data['monetary_values']:
                # Create line chart from monetary values
                data = statistical_data['monetary_values'][:6]
                labels = [item['context'] for item in data]
                values = [item['value'] for item in data]
                
                fig = px.line(x=labels, y=values, markers=True,
                            color_discrete_sequence=[self.colors[0]],
                            title="<b>Financial Metrics Trend</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
                
            elif chart_type == 'comparison_bar' and statistical_data['comparisons']:
                # Create comparison chart
                data = statistical_data['comparisons'][:6]
                labels = [item['context'] for item in data]
                values = [item['value'] if item['type'] == 'increase' else -item['value'] for item in data]
                colors = ['green' if v > 0 else 'red' for v in values]
                
                fig = px.bar(x=labels, y=values, color=colors,
                            title="<b>Performance Changes (%)</b>")
                fig.update_layout(**layout_config)
                fig.update_xaxis(tickangle=45)
                
            else:
                return ""
                
            # Convert to base64
            buffer = io.BytesIO()
            fig.write_image(buffer, format='png', scale=2)
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode()
            
        except Exception as e:
            logger.error(f"Error creating statistical chart: {str(e)}")
            return ""

    def suggest_chart_type(self, statistical_data: Dict[str, Any]) -> str:
        """Suggest the best chart type based on available data"""
        if statistical_data['percentages']:
            return 'percentage_pie'
        elif statistical_data['monetary_values']:
            return 'monetary_line'
        elif statistical_data['comparisons']:
            return 'comparison_bar'
        elif statistical_data['numbers']:
            return 'statistics_bar'
        return 'none'

class PresentationGenerator:
    def __init__(self, settings: Settings):
        self.settings = settings
        self.logger = logging.getLogger(__name__)
        self.chart_generator = ChartGenerator()
        self.default_theme = {
            "colors": {
                "primary": "#003087",
                "secondary": "#FFFFFF",
                "accent": "#FFA500",
                "text": "#404040",
                "footer": "#808080"
            },
            "fonts": {
                "title": {"name": "Calibri", "size": 44, "bold": True},
                "subtitle": {"name": "Calibri", "size": 24, "bold": False},
                "content": {"name": "Calibri", "size": 18, "bold": False},
                "footer": {"name": "Calibri", "size": 10, "bold": False}
            },
            "footer": {
                "text": "{company_name} | {date} | Confidential",
                "position": {"left": 0.5, "top": 7.0, "width": 9.0, "height": 0.5}
            }
        }
        self.theme = self.default_theme 
        try:
            theme_path = Path(__file__).parent / "templates" / "styles" / "corporate_theme.json"
            if theme_path.exists():
                with open(theme_path, "r") as f:
                    loaded_theme = json.load(f)
                    self.theme = {**self.default_theme, **loaded_theme.get("theme", {})}
                    for key in ["colors", "fonts", "footer"]:
                        if key in loaded_theme.get("theme", {}) and isinstance(loaded_theme["theme"][key], dict):
                            self.theme[key] = {**self.default_theme[key], **loaded_theme["theme"][key]}
            else:
                self.logger.warning(f"Theme file not found at {theme_path}. Using default theme.")
        except Exception as e:
            self.logger.warning(f"Failed to load theme: {str(e)}. Using default theme.")
        
        self.primary_color = RGBColor.from_string(self.theme["colors"]["primary"].lstrip("#"))
        self.secondary_color = RGBColor.from_string(self.theme["colors"]["secondary"].lstrip("#"))
        self.accent_color = RGBColor.from_string(self.theme["colors"]["accent"].lstrip("#"))
        self.text_color = RGBColor.from_string(self.theme["colors"]["text"].lstrip("#"))
        self.footer_color = RGBColor.from_string(self.theme["colors"]["footer"].lstrip("#"))
    
    def generate_presentation(self, content: str, company_name: str, 
                             report_month: str, report_year: int, 
                             domain_focus: List[str], language: str,
                             content_type: str, tables: List[Dict[str, List]], 
                             style: str, content_analysis: Dict[str, Any] = None) -> str:
        try:
            try:
                prs = Presentation(str(self.settings.template_path))
            except Exception as e:
                self.logger.warning(f"Template not found at {self.settings.template_path}: {str(e)}. Using default presentation.")
                st.warning(f"Using default presentation template as {self.settings.template_path} is missing.")
                prs = Presentation()
            
            bedrock_client = BedrockClient(self.settings)
            content_structure = bedrock_client.generate_presentation_content(
                content, company_name, report_month, report_year, domain_focus, language, content_type, tables, content_analysis
            )
            if not content_structure:
                content_structure = bedrock_client._create_default_structure(company_name, report_month, report_year, content_type, tables)
            
            self._create_enhanced_title_slide(prs, content_structure['title'], company_name, report_month, report_year, domain_focus, style)
            for slide_data in content_structure['slides']:
                self._create_enhanced_content_slide(prs, slide_data, style)
            
            # Add statistical slide if applicable
            if content_analysis and content_analysis.get('has_statistics', False):
                self._create_statistical_slide(prs, content_analysis['statistical_data'], style)
            
            filename = f"Presentation_{company_name.replace(' ', '_')}_{report_month}_{report_year}.pptx"
            output_path = self.settings.output_folder / filename
            prs.save(str(output_path))
            return str(output_path)
        except Exception as e:
            self.logger.error(f"Error generating presentation: {str(e)}")
            st.error(f"Error generating presentation: {str(e)}")
            raise
    
    def _create_enhanced_title_slide(self, prs: Presentation, title: str, company_name: str, 
                                    report_month: str, report_year: int, domain_focus: List[str], style: str):
        from pptx.enum.shapes import PP_PLACEHOLDER
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_placeholder = None
        subtitle_placeholder = None

        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_placeholder = shape
            elif shape.name == "Subtitle 2" or shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_placeholder = shape
        
        # Populate title
        if title_placeholder:
            text_frame = title_placeholder.text_frame
            text_frame.text = title
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
            paragraph.font.color.rgb = self.primary_color
            paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(2)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
            title_paragraph.font.color.rgb = self.primary_color
            title_paragraph.font.bold = self.theme["fonts"]["title"]["bold"]
            title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Populate subtitle
        subtitle_text = f"{company_name} | {report_month} {report_year} | {', '.join(domain_focus)}"
        if subtitle_placeholder:
            text_frame = subtitle_placeholder.text_frame
            text_frame.text = subtitle_text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
            paragraph.font.color.rgb = self.accent_color
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            subtitle_top = Inches(3.5)
            subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.font.size = Pt(self.theme["fonts"]["subtitle"]["size"])
            subtitle_paragraph.font.color.rgb = self.accent_color
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        if style in ["Creative", "Minimalist"]:
            self._add_decorative_image(slide, style)
        self._add_footer(slide, company_name, f"{report_month} {report_year}")
    
    def _create_enhanced_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], style: str):
        from pptx.enum.shapes import PP_PLACEHOLDER
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title_placeholder = None
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_placeholder = shape
            elif shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                body_placeholder = shape

        # Populate title
        if title_placeholder:
            title_shape = title_placeholder
            title_shape.text = slide_data['title']
            title_shape.text_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["title"]["size"])
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.primary_color
            title_shape.text_frame.paragraphs[0].font.bold = True
        else:
            left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
            title_paragraph.font.color.rgb = self.primary_color
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Populate content (bullet points or chart)
        if slide_data.get('chart_suggestion') != "none" and slide_data.get('chart_data'):
            content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            self._add_bullet_points(content_box.text_frame, slide_data['content'])

            chart_left = Inches(5.0)
            self._add_chart_to_slide(slide, slide_data['chart_suggestion'], slide_data['chart_data'], chart_left, content_top, style)
        else:
            if body_placeholder:
                content_box = body_placeholder
                content_box.text_frame.clear()
                self._add_bullet_points(content_box.text_frame, slide_data['content'])
            else:
                content_left, content_top, content_width, content_height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
                content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                self._add_bullet_points(content_box.text_frame, slide_data['content'])
        
        self._add_footer(slide, "Your Company", f"{datetime.now().strftime('%B %Y')}")
    
    def _create_statistical_slide(self, prs: Presentation, statistical_data: Dict[str, Any], style: str):
        """Create a slide specifically for statistical visualizations"""
        from pptx.enum.shapes import PP_PLACEHOLDER
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_placeholder = shape
                break
        
        if title_placeholder:
            title_placeholder.text = "📊 Statistical Analysis"
            title_placeholder.text_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["title"]["size"])
            title_placeholder.text_frame.paragraphs[0].font.color.rgb = self.primary_color
            title_placeholder.text_frame.paragraphs[0].font.bold = True
        else:
            left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = "📊 Statistical Analysis"
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(self.theme["fonts"]["title"]["size"])
            title_paragraph.font.color.rgb = self.primary_color
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
        
        # Add statistical charts
        chart_type = self.chart_generator.suggest_chart_type(statistical_data)
        if chart_type != 'none':
            chart_base64 = self.chart_generator.create_statistical_chart(statistical_data, chart_type, style)
            if chart_base64:
                buffer = io.BytesIO(base64.b64decode(chart_base64))
                slide.shapes.add_picture(buffer, Inches(1), Inches(2), Inches(8), Inches(4.5))
        
        # Add statistical summary text
        summary_text = self._generate_statistical_summary(statistical_data)
        if summary_text:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = summary_text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = self.text_color
        
        self._add_footer(slide, "Statistical Analysis", f"{datetime.now().strftime('%B %Y')}")
    
    def _generate_statistical_summary(self, statistical_data: Dict[str, Any]) -> str:
        """Generate a text summary of statistical findings"""
        summary_parts = []
        
        if statistical_data['numbers']:
            summary_parts.append(f"Key metrics: {len(statistical_data['numbers'])} data points identified")
        
        if statistical_data['percentages']:
            avg_percentage = sum(item['value'] for item in statistical_data['percentages']) / len(statistical_data['percentages'])
            summary_parts.append(f"Average percentage value: {avg_percentage:.1f}%")
        
        if statistical_data['monetary_values']:
            total_value = sum(item['value'] for item in statistical_data['monetary_values'])
            summary_parts.append(f"Total monetary value: ${total_value:,.0f}")
        
        return " | ".join(summary_parts)
    
    def _add_bullet_points(self, text_frame, content_list: List[str]):
        text_frame.clear()
        for i, item in enumerate(content_list):
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = item
            paragraph.font.size = Pt(self.theme["fonts"]["content"]["size"])
            paragraph.font.color.rgb = self.text_color
            paragraph.level = 0
            paragraph.space_after = Pt(12)
    
    def _add_chart_to_slide(self, slide, chart_type: str, chart_data: Dict[str, Any], left: float, top: float, style: str):
        try:
            chart_width, chart_height = Inches(4.5), Inches(3.5)
            if chart_type == 'table':
                table = slide.shapes.add_table(
                    len(chart_data.get('rows', [])[:5]) + 1,
                    len(chart_data.get('houses', [])),
                    left, top, chart_width, chart_height
                ).table
                for i, header in enumerate(chart_data.get('headers', [])):
                    table.cell(0, i).text = header
                    table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(14)
                for i, row in enumerate(chart_data.get('rows', [])[:5], 1):
                    for j, value in enumerate(row):
                        table.cell(i, j).text = str(value)
                        table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)
            else:
                chart_base64 = self.chart_generator.create_chart(chart_type, chart_data, style)
                if chart_base64:
                    buffer = io.BytesIO(base64.b64decode(chart_base64))
                    slide.shapes.add_picture(buffer, left, top, chart_width, chart_height)
        except Exception as e:
            self.logger.error(f"Error adding chart: {str(e)}")
            st.error(f"Error adding chart: {str(e)}")
            fallback_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(3))
            fallback_box.text_frame.text = "📊 Chart unavailable\n(Data visualization)"
            fallback_box.text_frame.paragraphs[0].font.size = Pt(16)
            fallback_box.text_frame.paragraphs[0].font.color.rgb = self.accent_color
    
    def _add_decorative_image(self, slide, style: str):
        """Add a decorative image or icon based on style"""
        try:
            img_filename = "creative_icon.png" if style == "Creative" else "minimal_icon.png"
            img_path = Path(__file__).parent / "templates" / "icons" / img_filename
            
            if not img_path.exists():
                img = Image.new('RGB', (200, 200), color=(self.accent_color.rgb[0], self.accent_color.rgb[1], self.accent_color.rgb[2]))
                img.save(img_path)

            left, top = Inches(8.5), Inches(0.5)
            width = Inches(1.5) if style == "Creative" else Inches(1)
            height = Inches(1.5) if style == "Creative" else Inches(1)

            slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)

        except Exception as e:
            self.logger.error(f"Error adding decorative image: {str(e)}")
            left, top, width, height = Inches(8.5), Inches(0.5), Inches(1), Inches(1)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.accent_color
            shape.line.fill.background()
    
    def _add_footer(self, slide, company_name: str, date: str):
        left = Inches(self.theme["footer"]["position"]["left"])
        top = Inches(self.theme["footer"]["position"]["top"])
        width = Inches(self.theme["footer"]["position"]["width"])
        height = Inches(self.theme["footer"]["position"]["height"])

        # Add shape for footer background
        footer_background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(10), height)
        footer_background.fill.solid()
        footer_background.fill.fore_color.rgb = self.secondary_color
        footer_background.line.fill.background()

        footer_box = slide.shapes.add_textbox(left, top, width, height)
        footer_frame = footer_box.text_frame
        footer_frame.text = self.theme["footer"]["text"].format(company_name=company_name, date=date)
        footer_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
        footer_frame.paragraphs[0].font.color.rgb = self.footer_color
        footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Add slide number
        slide_number_left = Inches(0.5)
        slide_number_box = slide.shapes.add_textbox(slide_number_left, top, Inches(1), height)
        slide_number_frame = slide_number_box.text_frame
        slide_number_frame.text = f"Page {slide.slide_id}"
        slide_number_frame.paragraphs[0].font.size = Pt(self.theme["fonts"]["footer"]["size"])
        slide_number_frame.paragraphs[0].font.color.rgb = self.footer_color
        slide_number_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

@st.cache_resource
def initialize_components():
    settings = Settings()
    doc_processor = DocumentProcessor()
    presentation_generator = PresentationGenerator(settings)
    return settings, doc_processor, presentation_generator

def main():
    st.markdown("""
    <div class="main-header">
        <h1>🎯 Smart Presentation Generator</h1>
        <p>Transform any content into beautiful, professional presentations with AI-powered insights and visualizations</p>
    </div>
    """, unsafe_allow_html=True)
    
    settings, doc_processor, presentation_generator = initialize_components()
    content_type = "Not yet processed"
    
    with st.sidebar:
        st.header("⚙️ Configuration")
        input_method = st.radio(
            "Choose input method:",
            ["📄 Upload PDF", "📝 Direct Text Input"],
            horizontal=True
        )
        target_language = st.selectbox(
            "Target Language",
            options=["en", "es", "fr", "de", "it", "pt", "nl", "ru", "zh", "ja", "ko", "ar", "hi"],
            format_func=lambda x: {
                "en": "🇺🇸 English", "es": "🇪🇸 Spanish", "fr": "🇫🇷 French", 
                "de": "🇩🇪 German", "it": "🇮🇹 Italian", "pt": "🇵🇹 Portuguese", 
                "nl": "🇳🇱 Dutch", "ru": "🇷🇺 Russian", "zh": "🇨🇳 Chinese", 
                "ja": "🇯🇵 Japanese", "ko": "🇰🇷 Korean", "ar": "🇸🇦 Arabic", 
                "hi": "🇮🇳 Hindi"
            }[x]
        )
        company_name = st.text_input("Company Name", "Your Company")
        col1, col2 = st.columns(2)
        with col1:
            report_month = st.selectbox(
                "Month",
                options=["January", "February", "March", "April", "May", "June",
                         "July", "August", "September", "October", "November", "December"],
                index=datetime.now().month - 1
            )
        with col2:
            report_year = st.number_input("Year", min_value=2020, max_value=2030, value=datetime.now().year)
        domain_focus = st.multiselect(
            "Focus Areas",
            options=["Logistics Vehicles", "Cargo Management", "Fleet Operations", 
                     "Supply Chain", "Transportation Efficiency", "Business Performance",
                     "Financial Reporting", "Market Analysis", "Strategic Planning",
                     "Narrative Development", "Character Analysis", "Plot Structure",
                     "Standard Operating Procedures", "Compliance and Regulations", "Operational Guidelines"],
            default=["Business Performance"]
        )
        presentation_style = st.selectbox(
            "Presentation Style",
            options=["Default", "Minimalist", "Creative"]
        )
        
        st.markdown("---")
        st.markdown("Developed by **Your Team Name**")

    document_content = ""
    tables_extracted: List[Dict[str, List]] = []
    content_analysis = None

    if input_method == "📄 Upload PDF":
        uploaded_file = st.file_uploader("Upload a PDF document", type=["pdf"])
        if uploaded_file:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                temp_pdf_path = tmp_file.name
            
            with st.spinner("Extracting text and analyzing content..."):
                document_content = doc_processor.extract_text_from_pdf(temp_pdf_path)
                tables_extracted = doc_processor.extract_tables(document_content)
                content_analysis = doc_processor.analyze_content_type(document_content)
                content_type = content_analysis['type']
                st.session_state['document_content'] = document_content
                st.session_state['tables_extracted'] = tables_extracted
                st.session_state['content_type'] = content_type
                st.session_state['content_analysis'] = content_analysis
            
            os.unlink(temp_pdf_path)
            if document_content:
                st.success("PDF processed successfully!")
                st.expander("Extracted Text Preview").write(document_content[:1000] + "...")
                if tables_extracted:
                    st.expander("Extracted Tables Preview").write(tables_extracted)
            else:
                st.error("Failed to extract content from PDF.")
    
    elif input_method == "📝 Direct Text Input":
        text_input = st.text_area("Paste your content here:", height=300)
        if text_input:
            with st.spinner("Analyzing content..."):
                document_content = doc_processor.process_text_document(text_input)
                tables_extracted = doc_processor.extract_tables(document_content)
                content_analysis = doc_processor.analyze_content_type(document_content)
                content_type = content_analysis['type']
                st.session_state['document_content'] = document_content
                st.session_state['tables_extracted'] = tables_extracted
                st.session_state['content_type'] = content_type
                st.session_state['content_analysis'] = content_analysis
            st.success("Text content processed successfully!")
    
    # Display statistical analysis in UI
    if 'content_analysis' in st.session_state and st.session_state['content_analysis'].get('has_statistics', False):
        st.subheader("📊 Statistical Analysis Found")
        
        statistical_data = st.session_state['content_analysis']['statistical_data']
        
        col1, col2 = st.columns(2)
        
        with col1:
            if statistical_data['numbers']:
                st.write("**Key Metrics:**")
                for item in statistical_data['numbers'][:5]:
                    st.write(f"• {item['context']}: {item['original']}")
        
        with col2:
            if statistical_data['percentages']:
                st.write("**Percentages:**")
                for item in statistical_data['percentages'][:5]:
                    st.write(f"• {item['context']}: {item['original']}")
        
        # Show chart preview
        chart_generator = ChartGenerator()
        chart_type = chart_generator.suggest_chart_type(statistical_data)
        if chart_type != 'none':
            chart_base64 = chart_generator.create_statistical_chart(statistical_data, chart_type, "Default")
            if chart_base64:
                st.image(base64.b64decode(chart_base64), caption="Statistical Chart Preview")

    if st.button("Generate Presentation 🚀", type="primary"):
        if 'document_content' not in st.session_state or not st.session_state['document_content']:
            st.error("Please upload a document or enter text content first.")
            return

        with st.spinner("Generating professional presentation... This may take a few moments."):
            try:
                # Retrieve processed content from session state
                doc_content = st.session_state['document_content']
                tables = st.session_state['tables_extracted']
                c_type = st.session_state['content_type']
                content_analysis = st.session_state.get('content_analysis', {})

                # Translate content if target language is not English
                if target_language != 'en':
                    bedrock_client = BedrockClient(settings)
                    doc_content = bedrock_client.translate_text(doc_content, target_language)
                
                output_ppt_path = presentation_generator.generate_presentation(
                    doc_content, company_name, report_month, report_year, domain_focus, target_language,
                    c_type, tables, presentation_style, content_analysis
                )
                
                if output_ppt_path and os.path.exists(output_ppt_path):
                    with open(output_ppt_path, "rb") as file:
                        btn = st.download_button(
                            label="Download Presentation ⬇️",
                            data=file,
                            file_name=os.path.basename(output_ppt_path),
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    st.success("Presentation generated successfully! Click the button above to download.")
                    st.balloons()
                else:
                    st.error("Failed to generate presentation. Please check logs for details.")
            except Exception as e:
                st.error(f"An unexpected error occurred during presentation generation: {e}")
                logger.error(f"Critical error in main generation flow: {e}")

    st.markdown("---")
    st.info("💡 **Tips for best results:**\n"
            "- Upload clear, well-formatted PDF documents.\n"
            "- Specify relevant focus areas for tailored content.\n"
            "- Ensure your AWS Bedrock credentials are set up in `.env` if you're using LLM functionality.")

if __name__ == "__main__":
    main()